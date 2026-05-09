#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate Lesson 2 teaching PPT from lesson2-spec.md."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors (same palette as lesson 1)
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
SUBTITLE_GRAY = RGBColor(0xA0, 0xA0, 0xB0)
DIM_GRAY = RGBColor(0x80, 0x80, 0x90)
ACCENT_BLUE = RGBColor(0x3A, 0x86, 0xFF)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
CODE_BG = RGBColor(0x0D, 0x11, 0x17)
YELLOW = RGBColor(0xF1, 0xC4, 0x0F)
RED_SOFT = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE = RGBColor(0xB5, 0x8B, 0xFF)


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_text(slide, left, top, width, height, text, size=18,
             color=LIGHT_GRAY, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.alignment = align
    return tf


def add_bullets(slide, left, top, width, height, items, size=16,
                color=LIGHT_GRAY, bullet_color=ORANGE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r1 = p.add_run()
        r1.text = '▸ '
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.name = 'Microsoft YaHei'
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = 'Microsoft YaHei'
    return tf


def add_code(slide, left, top, width, height, code, size=14, color=GREEN):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = 'Courier New'
    p.alignment = PP_ALIGN.LEFT
    return tf


def section_slide(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
             title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8),
                 subtitle, size=20, color=SUBTITLE_GRAY, align=PP_ALIGN.CENTER)
    return slide


def teaching_slide(time_tag, title, items_left, items_right=None,
                   note=None, code=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text(slide, Inches(0.8), Inches(0.3), Inches(2.5), Inches(0.4),
             time_tag, size=13, color=ORANGE, bold=True)
    add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
             title, size=28, color=WHITE, bold=True)

    if items_right:
        add_bullets(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(3.5),
                    items_left, size=16)
        add_bullets(slide, Inches(7), Inches(1.7), Inches(5.8), Inches(3.5),
                    items_right, size=16)
    else:
        add_bullets(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(3.5),
                    items_left, size=17)

    if code:
        add_code(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5), code)

    if note:
        add_text(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                 note, size=14, color=YELLOW, bold=True)
    return slide


# ============================================================
# Slide 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(0.6),
         'Claude Code 实战课程', size=20, color=SUBTITLE_GRAY,
         align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
         '第 2 次课', size=54, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
         '安装缓冲 + AI 思维模型 + 初次动手',
         size=28, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
         '今天的核心:学会"怎么跟 AI 对话",让它做出你真正想要的东西',
         size=16, color=DIM_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: Today's Goals
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
         '今天结束时,你会……', size=30, color=WHITE, bold=True)

goals = [
    ('✅', '全员安装就绪',
     '100% 学员完成本地 Claude Code 安装(或使用网页版备用方案)'),
    ('🧠', '掌握四个思维模型',
     '初级开发者 · 上下文 · 任务书 · 迭代'),
    ('🎨', '做出英语学习页面',
     '基于提供的 txt 文件做一个 Naval 金句学习页面,至少修改一次'),
    ('👀', '建立评估能力',
     '能判断 Claude 的输出对不对,并用自然语言提出修改要求'),
]
for i, (emoji, title, desc) in enumerate(goals):
    y = 1.5 + i * 1.3
    add_text(slide, Inches(0.8), Inches(y), Inches(0.8), Inches(0.6),
             emoji, size=28, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.8), Inches(y), Inches(4), Inches(0.5),
             title, size=22, color=ORANGE, bold=True)
    add_text(slide, Inches(1.8), Inches(y + 0.5), Inches(10.5), Inches(0.5),
             desc, size=16, color=LIGHT_GRAY)


# ============================================================
# Slide 3: Timeline Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
         '今天的时间线', size=30, color=WHITE, bold=True)

timeline = [
    ('0:00', '5 min', '开场回顾', '上次课做了什么,今天要学什么'),
    ('0:05', '15 min', '安装缓冲', '未完成的补装,完成的自由探索'),
    ('0:20', '10 min', '思维模型 1', 'Claude 是初级开发者,你是经理'),
    ('0:30', '10 min', '思维模型 2', '上下文决定一切(现场对比演示)'),
    ('0:40', '8 min', '思维模型 3', '先写任务书,再让 AI 动手'),
    ('0:48', '7 min', '思维模型 4', '迭代,而非追求完美'),
    ('0:55', '10 min', '休息', '☕'),
    ('1:05', '25 min', '练习 1', '基于 txt 文件做 Naval 金句英语学习页面'),
    ('1:30', '20 min', '练习 2', '修改页面:优化句子、补充音标翻译解读'),
    ('1:50', '10 min', '分享 + 收尾', '同伴展示 + 总结 + 预告下次课'),
]

for i, (time, dur, phase, desc) in enumerate(timeline):
    y = 1.3 + i * 0.55
    highlight = phase in ['思维模型 2', '练习 1', '练习 2']
    c = ORANGE if highlight else SUBTITLE_GRAY
    if phase == '休息':
        c = DIM_GRAY
    add_text(slide, Inches(0.8), Inches(y), Inches(0.9), Inches(0.4),
             time, size=14, color=DIM_GRAY, bold=True)
    add_text(slide, Inches(1.7), Inches(y), Inches(0.9), Inches(0.4),
             dur, size=12, color=DIM_GRAY)
    add_text(slide, Inches(2.8), Inches(y), Inches(3), Inches(0.4),
             phase, size=15, color=c, bold=True)
    add_text(slide, Inches(6.0), Inches(y), Inches(7), Inches(0.4),
             desc, size=14, color=LIGHT_GRAY)


# ============================================================
# Slide 4: Section - Opening
# ============================================================
section_slide('第一部分:开场回顾',
              '0:00 - 0:05 · 上次课做了什么,今天要学什么')


# ============================================================
# Slide 5: Opening Script
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:00 - 0:05', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '欢迎回来!上次课我们做了……',
         size=28, color=WHITE, bold=True)

last_lesson = [
    '看了一个 10 分钟做出网页的现场演示',
    '把 Claude Code 装到了你们的电脑上',
    '每个人都跑通了第一个网页',
]
add_bullets(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(2),
            last_lesson, size=20, spacing=Pt(10))

add_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.5),
         '今天要学一件更重要的事——\n怎么跟 AI 对话,才能让它做出你真正想要的东西。',
         size=24, color=ORANGE, bold=True)

add_text(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8),
         '会对话的人和不会对话的人,用同一个 AI,做出来的东西天差地别。\n今天之后,你就是那个"会对话的人"。',
         size=17, color=LIGHT_GRAY)


# ============================================================
# Slide 6: Section - Install Buffer
# ============================================================
section_slide('第二部分:安装缓冲',
              '0:05 - 0:20 · 两组并行,目标是全员就绪')


# ============================================================
# Slide 7: Install Buffer - Two Tracks
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:05 - 0:20', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '两组并行,15 分钟', size=28, color=WHITE, bold=True)

# Left: Group A - not installed
add_text(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.5),
         '🅰 未完成安装的同学', size=20, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(2.3), Inches(5.8), Inches(0.5),
         '跟助教一起完成,15 分钟足够',
         size=15, color=SUBTITLE_GRAY)
group_a = [
    'Node.js → Claude Code → API Key',
    '10 分钟仍未完成 → 切换 claude.ai/code 网页版',
    '底线:20 分钟时全班 100% 可用',
]
add_bullets(slide, Inches(0.8), Inches(3.0), Inches(5.8), Inches(2.5),
            group_a, size=15, bullet_color=ACCENT_BLUE)

# Right: Group B - already installed
add_text(slide, Inches(7), Inches(1.7), Inches(5.8), Inches(0.5),
         '🅱 已完成安装的同学', size=20, color=ORANGE, bold=True)
add_text(slide, Inches(7), Inches(2.3), Inches(5.8), Inches(0.5),
         '启动 Claude Code,试试下面 3 个挑战',
         size=15, color=SUBTITLE_GRAY)
group_b = [
    '帮我把上次做的网页加一个"关于我"的段落',
    '给网页加一个当前时间显示',
    '让网页背景变成渐变色',
]
add_bullets(slide, Inches(7), Inches(3.0), Inches(5.8), Inches(2.5),
            group_b, size=15, bullet_color=GREEN)

add_text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
         '检查点 (0:20):全班每个人现在都能启动 Claude Code 了吗?举手确认! 目标:100%',
         size=14, color=YELLOW, bold=True)


# ============================================================
# Slide 8: Section - Mental Models
# ============================================================
section_slide('第三部分:AI 思维模型',
              '0:20 - 0:55 · 这门课最值钱的东西')


# ============================================================
# Slide 9: 4 Models Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
         '接下来 35 分钟:四种思维方式',
         size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
         '不是某个命令,而是四种思维方式。掌握了,你用 AI 的效果会比别人好 10 倍。',
         size=16, color=SUBTITLE_GRAY)

models = [
    ('1', '初级开发者', '你是经理,Claude 是你新招的员工'),
    ('2', '上下文决定一切', '你给的信息越多、越具体,结果越好'),
    ('3', '先写任务书', '动手前花 2 分钟想清楚你要什么'),
    ('4', '迭代,而非完美', '先做出来,再改好——不是一次说对'),
]
for i, (num, title, desc) in enumerate(models):
    y = 2.3 + i * 1.1
    add_text(slide, Inches(1), Inches(y), Inches(0.8), Inches(0.7),
             num, size=44, color=ORANGE, bold=True)
    add_text(slide, Inches(2.2), Inches(y + 0.05), Inches(3.5), Inches(0.5),
             title, size=24, color=WHITE, bold=True)
    add_text(slide, Inches(6), Inches(y + 0.15), Inches(7), Inches(0.5),
             desc, size=17, color=LIGHT_GRAY)


# ============================================================
# Slide 10: Model 1 - Intro
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:20 - 0:30', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '思维模型 1:Claude 是初级开发者,你是经理',
         size=26, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.6),
         '想象你刚招了一个新员工。这个新员工有三个特点:',
         size=19, color=SUBTITLE_GRAY)

traits = [
    ('速度极快', '你让他做什么,几秒钟就做完了'),
    ('听话', '你说什么他就做什么'),
    ('需要你给方向', '你说"做个好看的东西"他就照做——但可能不是你想要的'),
]
for i, (t, d) in enumerate(traits):
    y = 2.8 + i * 0.9
    add_text(slide, Inches(1.2), Inches(y), Inches(3.5), Inches(0.5),
             f'第{["一","二","三"][i]},{t}', size=20, color=ORANGE, bold=True)
    add_text(slide, Inches(5.5), Inches(y + 0.05), Inches(7.5), Inches(0.5),
             d, size=16, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8),
         'Claude Code 就是这样一个"初级开发者"。而你,是它的经理。',
         size=20, color=WHITE, bold=True)


# ============================================================
# Slide 11: Model 1 - Manager's 3 Jobs
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '思维模型 1', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '经理做什么?三件事',
         size=28, color=WHITE, bold=True)

jobs = [
    ('🎯', '给方向', '告诉它做什么'),
    ('🔍', '审查', '看看做得对不对'),
    ('💬', '反馈', '告诉它哪里要改'),
]
for i, (emoji, t, d) in enumerate(jobs):
    y = 2.0 + i * 1.2
    add_text(slide, Inches(1.5), Inches(y), Inches(0.8), Inches(0.8),
             emoji, size=36, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.8), Inches(y + 0.1), Inches(3), Inches(0.5),
             t, size=26, color=ORANGE, bold=True)
    add_text(slide, Inches(6.5), Inches(y + 0.2), Inches(6), Inches(0.5),
             d, size=20, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1),
         '你不需要知道代码怎么写。\n你只需要知道"我想要什么"和"这个结果对不对"。',
         size=20, color=YELLOW, bold=True)


# ============================================================
# Slide 12: Model 1 - PPT Analogy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '思维模型 1 · 类比', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '假设你让新员工做一份 PPT,你会怎么说?',
         size=26, color=WHITE, bold=True)

# Option A - bad
add_text(slide, Inches(0.8), Inches(2), Inches(5.8), Inches(0.5),
         '❌ A 同事', size=22, color=RED_SOFT, bold=True)
add_text(slide, Inches(0.8), Inches(2.6), Inches(5.8), Inches(1),
         '"帮我做个 PPT。"',
         size=22, color=WHITE)
add_text(slide, Inches(0.8), Inches(4.0), Inches(5.8), Inches(1.5),
         '结果:他做出来大概率不是你想要的。',
         size=16, color=SUBTITLE_GRAY)

# Option B - good
add_text(slide, Inches(7), Inches(2), Inches(5.8), Inches(0.5),
         '✅ B 同事', size=22, color=GREEN, bold=True)
add_text(slide, Inches(7), Inches(2.6), Inches(5.8), Inches(1.7),
         '"帮我做个 10 页的 PPT,主题是年度总结,第一页放标题和日期,第二页放今年的三个亮点,第三页放数据图表……"',
         size=16, color=WHITE)
add_text(slide, Inches(7), Inches(5.0), Inches(5.8), Inches(1),
         '结果:好多了,对吧?',
         size=16, color=SUBTITLE_GRAY)

add_text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
         '跟 Claude Code 对话也是一样的道理。你说得越清楚,它做得越对。',
         size=17, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 13: Model 2 - Context is Everything
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:30 - 0:40', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '思维模型 2:上下文决定一切',
         size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1),
         'Claude 只知道你告诉它的事情。\n它看不到你脑子里的画面。',
         size=28, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.8),
         '让我演示一下区别有多大——',
         size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
         '下一页:现场对比演示(本环节的重头戏)',
         size=14, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 14: Model 2 - Live Demo Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '现场对比演示', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '同一个 AI,不同的上下文',
         size=28, color=WHITE, bold=True)

# Left column - bad
add_text(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(0.5),
         '❌ 差的上下文', size=20, color=RED_SOFT, bold=True)
add_code(slide, Inches(0.8), Inches(2.2), Inches(5.8), Inches(0.9),
         'mkdir demo-bad && cd demo-bad\nclaude', size=13)
add_code(slide, Inches(0.8), Inches(3.4), Inches(5.8), Inches(0.7),
         '做一个网页', size=14, color=WHITE)
add_text(slide, Inches(0.8), Inches(4.5), Inches(5.8), Inches(2),
         '→ 能用吗?能。\n→ 但它不知道你想要什么,\n  所以做了一个非常普通的东西。',
         size=14, color=SUBTITLE_GRAY)

# Right column - good
add_text(slide, Inches(7), Inches(1.6), Inches(5.8), Inches(0.5),
         '✅ 好的上下文', size=20, color=GREEN, bold=True)
add_code(slide, Inches(7), Inches(2.2), Inches(5.8), Inches(0.9),
         'mkdir demo-good && cd demo-good\nclaude', size=13)
good_prompt = (
    '帮我做一个个人介绍网页。我叫小明,\n'
    '是一个摄影爱好者。\n'
    '- 顶部大标题"小明的世界"\n'
    '- 下面 3-4 句自我介绍\n'
    '- 三个板块:作品/爱好/联系方式\n'
    '- 文艺风,暖色调,字体柔和\n'
    '- 加一些好看的 emoji 装饰'
)
add_code(slide, Inches(7), Inches(3.2), Inches(5.8), Inches(3.2),
         good_prompt, size=12, color=WHITE)

add_text(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
         '"看到区别了吗?同一个 AI,结果天差地别。唯一的区别就是——你给了它多少信息。"',
         size=13, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 15: Model 2 - The Rule
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '思维模型 2 · 记住', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '记住这条规则',
         size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5),
         '你给 Claude 的信息越多、越具体,\n它给你的结果越好。',
         size=32, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.6),
         '不是要你写一篇文章。而是想清楚三个问题:',
         size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

qs = [
    ('我想要什么?', '功能、内容'),
    ('我想要什么样的?', '风格、感觉'),
    ('有没有具体要求?', '细节、限制'),
]
for i, (q, a) in enumerate(qs):
    x = 0.8 + i * 4.2
    add_text(slide, Inches(x), Inches(5.4), Inches(4), Inches(0.5),
             q, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(x), Inches(6.0), Inches(4), Inches(0.5),
             a, size=16, color=SUBTITLE_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 16: Model 3 - Write Spec First
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:40 - 0:48', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '思维模型 3:先写任务书,再让 AI 动手',
         size=26, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1),
         '动手让 AI 做之前,先花 2 分钟把你想要的东西写下来。',
         size=22, color=ORANGE, bold=True)

add_text(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.6),
         '就像你给新员工派活,先把需求写在纸上,而不是嘴上一句"你看着办"。',
         size=17, color=LIGHT_GRAY)

# Task book template
add_text(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.5),
         '任务书三件套', size=20, color=WHITE, bold=True)

template = [
    ('1.', '我想要什么?', '功能、内容'),
    ('2.', '我想要什么样的?', '风格、感觉'),
    ('3.', '有没有具体要求?', '细节、限制'),
]
for i, (num, q, hint) in enumerate(template):
    y = 4.7 + i * 0.7
    add_text(slide, Inches(1), Inches(y), Inches(0.5), Inches(0.4),
             num, size=20, color=ORANGE, bold=True)
    add_text(slide, Inches(1.7), Inches(y), Inches(5), Inches(0.4),
             q, size=20, color=WHITE, bold=True)
    add_text(slide, Inches(8), Inches(y + 0.05), Inches(5), Inches(0.4),
             f'— {hint}', size=17, color=SUBTITLE_GRAY)


# ============================================================
# Slide 17: Model 3 - Why Write First
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '思维模型 3 · 为什么', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '为什么要先写?三个原因',
         size=28, color=WHITE, bold=True)

reasons = [
    ('1', '逼你自己想清楚',
     '很多时候你以为知道自己想要什么,写下来才发现"其实我没想好"'),
    ('2', 'AI 一次就能做对大半',
     '省掉反复返工的时间'),
    ('3', '任务书可以复用',
     '下次想做类似的东西,改改任务书就行'),
]
for i, (num, title, desc) in enumerate(reasons):
    y = 1.8 + i * 1.2
    add_text(slide, Inches(1.2), Inches(y), Inches(0.8), Inches(0.7),
             num, size=44, color=ORANGE, bold=True)
    add_text(slide, Inches(2.5), Inches(y + 0.05), Inches(10), Inches(0.5),
             title, size=22, color=WHITE, bold=True)
    add_text(slide, Inches(2.5), Inches(y + 0.6), Inches(10), Inches(0.5),
             desc, size=15, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8),
         '💡 这是程序员现在用 AI 的标准做法(spec-driven),\n你们一上来就学的是高级打法。',
         size=16, color=PURPLE, bold=True)


# ============================================================
# Slide 18: Model 3 - Retrospective on Xiaoming
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '思维模型 3 · 回溯', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '其实刚才"小明的世界"就是一份任务书',
         size=24, color=WHITE, bold=True)

items = [
    ('我想要什么?', '个人介绍页面,主人公是摄影爱好者小明'),
    ('我想要什么样的?', '文艺风格、暖色调、字体柔和'),
    ('有没有具体要求?', '大标题"小明的世界"、三个板块、加 emoji'),
]
for i, (q, a) in enumerate(items):
    y = 2.0 + i * 1.3
    add_text(slide, Inches(1), Inches(y), Inches(4.5), Inches(0.5),
             q, size=20, color=ORANGE, bold=True)
    add_text(slide, Inches(1), Inches(y + 0.55), Inches(11.5), Inches(0.6),
             f'→ {a}', size=18, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.6),
         '三个问题回答清楚,任务书就成了。不复杂。',
         size=18, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 19: Model 4 - Iterate
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:48 - 0:55', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '思维模型 4:迭代,而非追求完美',
         size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1),
         '误区:觉得要一次把需求说完美,AI 才能做出好东西。',
         size=19, color=RED_SOFT)

add_text(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.6),
         '事实上,最好的用法是:',
         size=18, color=SUBTITLE_GRAY)

steps = [
    ('1', '先说个大概', '→ 看结果'),
    ('2', '不满意的地方', '→ 告诉它改'),
    ('3', '再看', '→ 再改'),
    ('4', '几轮之后', '→ 就是你想要的了'),
]
for i, (n, a, b) in enumerate(steps):
    y = 3.6 + i * 0.6
    add_text(slide, Inches(1.5), Inches(y), Inches(0.6), Inches(0.4),
             n + '.', size=22, color=ORANGE, bold=True)
    add_text(slide, Inches(2.3), Inches(y + 0.05), Inches(4), Inches(0.4),
             a, size=19, color=WHITE, bold=True)
    add_text(slide, Inches(6.5), Inches(y + 0.1), Inches(5), Inches(0.4),
             b, size=17, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
         '这就像改作文。没有人第一稿就是完美的。',
         size=17, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 20: Model 4 - Live Iteration Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '现场迭代演示', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '用刚才的小明页面,继续改',
         size=28, color=WHITE, bold=True)

# Round 1
add_text(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(0.5),
         '第 1 轮修改', size=20, color=ACCENT_BLUE, bold=True)
round1 = (
    '这个页面不错,但我想改几个地方:\n'
    '1. 标题太小了,再大一倍\n'
    '2. 背景色太亮了,换成米白色\n'
    '3. "联系方式"板块加一个微信二维码占位图'
)
add_code(slide, Inches(0.8), Inches(2.2), Inches(5.8), Inches(2.5),
         round1, size=13, color=WHITE)

# Round 2
add_text(slide, Inches(7), Inches(1.6), Inches(5.8), Inches(0.5),
         '第 2 轮修改', size=20, color=ACCENT_BLUE, bold=True)
round2 = (
    '很好!再改一下:\n'
    '- 三个板块用卡片式布局,加圆角和阴影\n'
    '- 页面底部加一行小字:"用 AI 制作"'
)
add_code(slide, Inches(7), Inches(2.2), Inches(5.8), Inches(2.5),
         round2, size=13, color=WHITE)

add_text(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
         '我从来没有一次说完所有需求。我是一步步改过来的。',
         size=18, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.2),
         '核心技能:不是一次说对,而是一步步调对。',
         size=22, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 21: Model 4 - Three Rules
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '思维模型 4 · 记住', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '迭代三句话',
         size=28, color=WHITE, bold=True)

rules = [
    ('1', '第一次结果是草稿,不是成品',
     '别期望一次完美'),
    ('2', '改比重做快',
     '说"把标题改大"比重新描述整个页面快得多'),
    ('3', '每次只改 1-2 个地方',
     '不要一口气提 10 个修改'),
]
for i, (n, title, desc) in enumerate(rules):
    y = 2.0 + i * 1.4
    add_text(slide, Inches(1.2), Inches(y), Inches(0.8), Inches(0.7),
             n, size=44, color=ORANGE, bold=True)
    add_text(slide, Inches(2.5), Inches(y + 0.05), Inches(10), Inches(0.6),
             title, size=22, color=WHITE, bold=True)
    add_text(slide, Inches(2.5), Inches(y + 0.7), Inches(10), Inches(0.5),
             desc, size=16, color=LIGHT_GRAY)


# ============================================================
# Slide 22: Break
# ============================================================
section_slide('☕ 休息 10 分钟',
              '有疑问?休息时找老师或助教 · 回来后你们自己动手做')


# ============================================================
# Slide 23: Section - Practice
# ============================================================
section_slide('第四部分:引导式练习',
              '1:05 - 1:50 · 用四个思维模型做一个英语学习小项目')


# ============================================================
# Slide 24: Exercise 1 - Write Task Book First
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '1:05 - 1:30 · 练习 1', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '基于 naval.txt 做一个 Naval 金句英语学习页面',
         size=24, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.6),
         '动手敲命令之前,先做一件事 —— 写任务书',
         size=20, color=ORANGE, bold=True)

add_text(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.5),
         '老师会提供 naval.txt。先花 3 分钟回答三个问题:',
         size=17, color=SUBTITLE_GRAY)

q_items = [
    ('我想要什么?',
     '从 naval.txt 里挑 5 句话,做成英语学习页面'),
    ('我想要什么样的?',
     '简约?卡片式?像背单词 App?什么颜色?'),
    ('有没有具体要求?',
     '是否加关键词高亮?是否标注来源?是否加重点词解释?'),
]
for i, (q, hint) in enumerate(q_items):
    y = 3.2 + i * 1.0
    add_text(slide, Inches(1), Inches(y), Inches(5), Inches(0.5),
             f'{i+1}. {q}', size=19, color=WHITE, bold=True)
    add_text(slide, Inches(1.5), Inches(y + 0.55), Inches(11), Inches(0.5),
             hint, size=15, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
         '写完再动手。这 3 分钟花得值。',
         size=18, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 25: Exercise 1 - Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '练习 1 · 操作步骤', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '7 个步骤,跟着做',
         size=28, color=WHITE, bold=True)

step_items = [
    ('1', '写任务书', '(纸上或记事本) 回答三个问题'),
    ('2', '打开终端', 'macOS: Command+空格,输入 Terminal'),
    ('3', '创建项目文件夹', 'mkdir naval-quotes && cd naval-quotes'),
    ('4', '准备 txt 文件', '确保老师提供的 naval.txt 在当前文件夹'),
    ('5', '启动 Claude Code', 'claude'),
    ('6', '把任务书告诉 Claude', '请基于 naval.txt 选 5 句话,展示原句、音标、翻译和解读'),
    ('7', '打开浏览器查看', 'open index.html  /  start index.html'),
]
for i, (num, title, cmd) in enumerate(step_items):
    y = 1.45 + i * 0.7
    add_text(slide, Inches(1), Inches(y), Inches(0.6), Inches(0.5),
             num, size=24, color=ORANGE, bold=True)
    add_text(slide, Inches(2), Inches(y + 0.05), Inches(3.8), Inches(0.5),
             title, size=18, color=WHITE, bold=True)
    add_code(slide, Inches(5.8), Inches(y + 0.1), Inches(7), Inches(0.5),
             cmd, size=13)

add_text(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
         '检查点 (1:30):已经看到 Naval 金句学习页面的同学举手! 目标:90%+',
         size=14, color=YELLOW, bold=True)


# ============================================================
# Slide 26: Teacher's Role in Exercise 1
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2.5), Inches(0.4),
         '练习 1 · 老师/助教备注', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '老师和助教怎么做(学员看不到)',
         size=26, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.4),
         '(此页不展示给学员,仅老师参考)',
         size=13, color=DIM_GRAY)

# Teacher
add_text(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(0.5),
         '👨‍🏫 老师', size=20, color=ACCENT_BLUE, bold=True)
teacher = [
    '前 3 分钟:走动看学员的任务书,不看电脑',
    '太模糊 → 引导加细节:"5 句话要偏短还是偏有启发性?"',
    '太复杂 → 引导简化:"先做简单版"',
    '不要替学员写任务书,只提问引导',
    '3 分钟后放行,开始敲命令',
]
add_bullets(slide, Inches(0.8), Inches(2.6), Inches(5.8), Inches(4),
            teacher, size=14)

# TA
add_text(slide, Inches(7), Inches(2.0), Inches(5.8), Inches(0.5),
         '🧑‍💻 助教', size=20, color=GREEN, bold=True)
ta = [
    '帮助遇到技术问题的学员',
    '(终端操作、文件路径等)',
    '观察任务书质量',
    '常见困惑课后反馈给老师',
    '未成功学员:通常是路径或 Claude Code 未启动',
]
add_bullets(slide, Inches(7), Inches(2.6), Inches(5.8), Inches(4),
            ta, size=14, bullet_color=GREEN)


# ============================================================
# Slide 27: Exercise 2 - Modify Your Page
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '1:30 - 1:50 · 练习 2', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '修改你的 Naval 学习页面',
         size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.6),
         '看看你的页面,有没有什么地方不满意?',
         size=20, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.5),
         '句子太难?音标不统一?翻译太生硬?解读太浅?',
         size=17, color=SUBTITLE_GRAY)

add_text(slide, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.6),
         '用刚学的思维模型 4 —— 迭代。告诉 Claude 你想改什么。',
         size=20, color=ORANGE, bold=True)

add_text(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.6),
         '🎯 目标:至少改两次。每次只改 1-2 个地方。',
         size=22, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8),
         '下一页:修改示例参考(不要求照做,只是给你灵感)',
         size=14, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 28: Exercise 2 - Examples Table
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '练习 2 · 修改示例', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '你觉得…… 就这样跟 Claude 说',
         size=26, color=WHITE, bold=True)

# Header
add_text(slide, Inches(0.8), Inches(1.6), Inches(3.5), Inches(0.4),
         '你觉得……', size=16, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(4.5), Inches(1.6), Inches(8.3), Inches(0.4),
         '你可以这样跟 Claude 说', size=16, color=ACCENT_BLUE, bold=True)

examples = [
    ('句子太难', '把 5 句话换成更短、更适合英语初学者背诵的句子'),
    ('想加内容', '每句话下面再补一行关键词解释,\n告诉我最值得学的 1-2 个词'),
    ('布局太挤', '每张句子卡片之间加大间距,\n原文、音标、翻译、解读分层显示'),
    ('想加效果', '点击卡片后再展开中文翻译和解读,\n先默认只显示英文原句'),
    ('不够像学习工具', '这个页面还不像英语学习产品,帮我美化一下,\n参考简洁的背单词 App 风格'),
]
for i, (feel, say) in enumerate(examples):
    y = 2.2 + i * 0.85
    add_text(slide, Inches(0.8), Inches(y), Inches(3.5), Inches(0.5),
             feel, size=15, color=WHITE, bold=True)
    add_text(slide, Inches(4.5), Inches(y), Inches(8.3), Inches(0.8),
             say, size=14, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.5),
         '说不清也没关系,就用大白话描述感受,Claude 会帮你想办法。',
         size=14, color=YELLOW, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 29: Section - Share + Wrap
# ============================================================
section_slide('第五部分:分享 + 收尾',
              '1:50 - 2:00 · 同伴展示 + 总结 + 预告下次课')


# ============================================================
# Slide 30: Peer Share
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '1:50 - 1:55 · 同伴展示', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '转向你旁边的同学',
         size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2), Inches(11.5), Inches(1),
         '把你的 Naval 学习页面给 ta 看。',
         size=30, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(1),
         '互相看看,说两句:\n"你这个怎么做的?"\n"你跟 Claude 说了什么?"',
         size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6), Inches(11.5), Inches(0.6),
         '⏱ 给你们 3 分钟',
         size=22, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 31: Summary - 4 Mental Models
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '1:55 - 1:58 · 总结', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '今天学了四个思维模型',
         size=28, color=WHITE, bold=True)

summary = [
    ('1', '你是经理,Claude 是你的初级开发者',
     '你给方向、审查结果、提反馈'),
    ('2', '上下文决定一切',
     '你说得越清楚,它做得越好'),
    ('3', '先写任务书,再让 AI 动手',
     '动手之前花 2-3 分钟想清楚你要什么'),
    ('4', '迭代,不要追求一次完美',
     '先做出来,再改好'),
]
for i, (num, title, desc) in enumerate(summary):
    y = 1.6 + i * 1.05
    add_text(slide, Inches(1), Inches(y), Inches(0.7), Inches(0.6),
             num, size=36, color=ORANGE, bold=True)
    add_text(slide, Inches(2), Inches(y + 0.05), Inches(10.5), Inches(0.5),
             title, size=20, color=WHITE, bold=True)
    add_text(slide, Inches(2), Inches(y + 0.6), Inches(10.5), Inches(0.4),
             desc, size=15, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.6),
         '而且你们刚刚做到了 —— 每个人都写了任务书、基于 txt 做出了页面、改了至少两次。',
         size=15, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 32: Homework + Preview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '1:58 - 2:00', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7),
         '课后任务(可选,不强制)',
         size=28, color=WHITE, bold=True)

add_bullets(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(2.5),
    [
        '继续美化你的 Naval 金句学习页面 — 加关键词解释、收藏功能、换风格',
        '试试把材料换成另一份英文 txt 文件,做一个新的英语学习页面',
        '做不出来不用焦虑 — 记下你想做什么,下次课我帮你',
    ], size=17)

add_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.5),
         '做了的同学可以截图发微信群,我们互相看看!',
         size=15, color=YELLOW)

add_text(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.7),
         '下次课预告', size=24, color=ORANGE, bold=True)

add_text(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2),
         '第 3 次课:我们开始做一个更大的东西 —— 一个真正的小工具。\n我会带着你们一步一步做。',
         size=18, color=LIGHT_GRAY)


# ============================================================
# Slide 33: Teacher Notes - Backup Plans
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
         '老师备注:备用方案',
         size=28, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4),
         '(此页不展示给学员,仅老师参考)',
         size=13, color=DIM_GRAY)

backups = [
    ('安装缓冲超时', '立即切换 claude.ai/code 网页版,不再纠结本地安装'),
    ('对比演示效果不明显', '用预先准备的截图对比替代(提前测试并截好图)'),
    ('API 速率限制', '分批启动(先左半区,1 分钟后右半区);或切换备用 Key'),
    ('学员写不出提示词', '提供模板:"请基于 naval.txt 选 5 句话,页面风格___,每句展示原句/音标/翻译/解读"'),
    ('学员缺第 1 次课', '助教一对一帮助 + 发第 1 次课回顾讲义'),
    ('练习时间不够', '砍思维模型 3 到 5 分钟,或压缩练习 2 为 10 分钟'),
]
for i, (scenario, solution) in enumerate(backups):
    y = 1.8 + i * 0.85
    add_text(slide, Inches(0.8), Inches(y), Inches(3.5), Inches(0.4),
             scenario, size=16, color=RED_SOFT, bold=True)
    add_text(slide, Inches(4.5), Inches(y), Inches(8.3), Inches(0.5),
             solution, size=15, color=LIGHT_GRAY)


# ============================================================
# Slide 34: Teacher Notes - Key Reminders
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
         '老师备注:6 条关键提醒',
         size=28, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4),
         '(此页不展示给学员,仅老师参考)',
         size=13, color=DIM_GRAY)

reminders = [
    ('对比演示是本课的灵魂',
     '差的上下文 vs 好的上下文,差异必须让学员"哇"出来。提前测试至少 3 次。'),
    ('用类比,不要用术语',
     '经理管下属、改作文、装修房子——用学员生活中的事情解释。'),
    ('练习时间是神圣的',
     '思维模型讲解超时,砍讲解而不是砍练习。动手记忆是听讲的 10 倍。'),
    ('不要替学员写提示词',
     '"你想要什么风格?"比"你应该写简约风格"更好。'),
    ('第 2 次课是信心关键课',
     '目标是把第 1 次课的怀疑转化为"我确实能做到"。每人都要带着作品离开。'),
    ('同伴展示优于全班展示',
     '这阶段作品质量差异大,同伴展示更安全。全班展示留给第 6、8 次课。'),
]
for i, (title, desc) in enumerate(reminders):
    y = 1.7 + i * 0.82
    add_text(slide, Inches(0.8), Inches(y), Inches(0.5), Inches(0.4),
             str(i + 1), size=18, color=ORANGE, bold=True)
    add_text(slide, Inches(1.3), Inches(y), Inches(4), Inches(0.4),
             title, size=15, color=WHITE, bold=True)
    add_text(slide, Inches(5.2), Inches(y), Inches(7.8), Inches(0.4),
             desc, size=13, color=LIGHT_GRAY)


# ============================================================
# Save
# ============================================================
output = '/Users/wenfengzhu/ClaudeProjects/ClaudeClass/第2次课-教学PPT.pptx'
prs.save(output)
print(f'PPT saved to: {output}')
print(f'Total slides: {len(prs.slides)}')
