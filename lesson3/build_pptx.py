#!/usr/bin/env python3
"""Generate Apple-style .pptx for Harness Engineering (30 slides, Chinese)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ---------- Color palette (Apple-style) ----------
INK         = RGBColor(0x1d, 0x1d, 0x1f)
INK_SOFT    = RGBColor(0x6e, 0x6e, 0x73)
PAPER       = RGBColor(0xfb, 0xfb, 0xfd)
BLACK       = RGBColor(0x00, 0x00, 0x00)
WHITE       = RGBColor(0xff, 0xff, 0xff)
GRAY_LIGHT  = RGBColor(0xf5, 0xf5, 0xf7)
GRAY_TEXT   = RGBColor(0xa1, 0xa1, 0xa6)
LINE        = RGBColor(0xd2, 0xd2, 0xd7)
ACCENT      = RGBColor(0x00, 0x71, 0xe3)
ACCENT_SOFT = RGBColor(0xe6, 0xf2, 0xff)
WARM        = RGBColor(0xff, 0x9f, 0x0a)
WARM_SOFT   = RGBColor(0xff, 0xf7, 0xed)
GREEN       = RGBColor(0x34, 0xc7, 0x59)
GREEN_SOFT  = RGBColor(0xf0, 0xfd, 0xf4)
RED         = RGBColor(0xff, 0x3b, 0x30)
RED_SOFT    = RGBColor(0xfe, 0xf2, 0xf2)
PURPLE      = RGBColor(0xaf, 0x52, 0xde)
TINT_BLUE   = RGBColor(0xf0, 0xf7, 0xff)
TINT_WARM   = RGBColor(0xff, 0xfb, 0xeb)
DARK_CARD   = RGBColor(0x1c, 0x1c, 0x1e)
CALLOUT_BG  = RGBColor(0xff, 0xfb, 0xeb)
CALLOUT_BR  = RGBColor(0xb4, 0x53, 0x09)

# ---------- Setup 16:9 widescreen ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

FONT_HAN = "PingFang SC"
FONT_LATIN = "SF Pro Display"
FONT_MONO = "SF Mono"


# ---------- Helpers ----------
def add_slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    bg_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.shadow.inherit = False
    return s


def set_run(run, text, *, size=18, bold=False, color=INK, font_latin=FONT_LATIN, font_han=FONT_HAN):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_latin
    rPr = run._r.get_or_add_rPr()
    # Set East Asian font
    for tag in ("ea", "cs"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
    ea = rPr.makeelement(qn("a:ea"), {"typeface": font_han})
    rPr.append(ea)
    cs = rPr.makeelement(qn("a:cs"), {"typeface": font_han})
    rPr.append(cs)


def add_text(slide, x, y, w, h, lines,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             fill=None, font_latin=FONT_LATIN):
    """lines: list of (text, dict(size, bold, color))."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = anchor

    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if "space_before" in opts:
            p.space_before = Pt(opts["space_before"])
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])
        if "line_spacing" in opts:
            p.line_spacing = opts["line_spacing"]
        run = p.add_run()
        set_run(run,
                text,
                size=opts.get("size", 18),
                bold=opts.get("bold", False),
                color=opts.get("color", INK),
                font_latin=opts.get("font_latin", font_latin))
    return tb


def add_rect(slide, x, y, w, h, fill, line=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        # adjust corner radius
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, weight=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def slide_num(slide, n, total=30, dark=False):
    add_text(slide, Inches(11.5), Inches(0.3), Inches(1.5), Inches(0.3),
             [(f"{n:02d} · {total}", {"size": 11, "color": GRAY_TEXT if dark else INK_SOFT})],
             align=PP_ALIGN.RIGHT)


def eyebrow(slide, x, y, text, dark=False):
    add_text(slide, x, y, Inches(8), Inches(0.4),
             [(text.upper(), {"size": 12, "bold": True,
                              "color": GRAY_TEXT if dark else INK_SOFT})])


def callout(slide, x, y, w, h, text_lines):
    add_rect(slide, x, y, w, h, CALLOUT_BG, radius=0.06)
    # Left orange bar
    add_rect(slide, x, y, Inches(0.06), h, WARM)
    add_text(slide, x + Inches(0.25), y + Inches(0.15), w - Inches(0.4), h - Inches(0.3),
             text_lines, anchor=MSO_ANCHOR.MIDDLE)


def card(slide, x, y, w, h, content, *, dark=False, bg=None, accent_color=None):
    """content: list of (label, title, paragraphs[]) or simpler structures.
    Render an Apple-style card."""
    fill = bg if bg else (DARK_CARD if dark else WHITE)
    add_rect(slide, x, y, w, h, fill,
             line=LINE if not dark else DARK_CARD, radius=0.10)
    pad = Inches(0.3)
    inner_x = x + pad
    inner_y = y + pad
    inner_w = w - pad * 2
    label, title, paragraphs = content
    title_color = accent_color if accent_color else (WHITE if dark else INK)
    body_color  = GRAY_TEXT if dark else INK_SOFT
    label_color = GRAY_TEXT if dark else INK_SOFT

    lines = []
    if label:
        lines.append((label.upper(), {"size": 10, "bold": True,
                                      "color": label_color, "space_after": 4}))
    if title:
        lines.append((title, {"size": 18, "bold": True,
                              "color": title_color, "space_after": 6}))
    for p in paragraphs:
        lines.append((p, {"size": 12, "color": body_color,
                          "line_spacing": 1.35, "space_after": 3}))
    add_text(slide, inner_x, inner_y, inner_w, h - pad * 2, lines)


def gradient_title(slide, x, y, w, h, text):
    """Approximate gradient by rendering the text in a strong purple-pink color."""
    add_text(slide, x, y, w, h,
             [(text, {"size": 80, "bold": True, "color": PURPLE})])


# ============================================================
#                     SLIDES
# ============================================================

# ---------- 1. Cover ----------
s = add_slide(BLACK)
eyebrow(s, Inches(0.9), Inches(0.8), "A Beginner's Guide · 2026", dark=True)
add_text(s, Inches(0.9), Inches(1.4), Inches(11), Inches(2.0),
         [("Harness", {"size": 88, "bold": True, "color": WHITE})])
add_text(s, Inches(0.9), Inches(2.7), Inches(11), Inches(2.0),
         [("Engineering", {"size": 88, "bold": True, "color": PURPLE})])
add_text(s, Inches(0.9), Inches(4.7), Inches(11), Inches(1.5),
         [("给 AI 编程助手装上"缰绳" ——", {"size": 22, "color": GRAY_TEXT}),
          ("让它写出真正能用的代码。",   {"size": 22, "color": GRAY_TEXT})])
add_text(s, Inches(0.9), Inches(6.7), Inches(11), Inches(0.4),
         [("Based on Martin Fowler · Anthropic Engineering · 30 页 · 为零基础成年人整理",
           {"size": 11, "color": GRAY_TEXT})])

# ---------- 2. Why this matters ----------
s = add_slide(TINT_BLUE)
eyebrow(s, Inches(0.9), Inches(0.6), "为什么你该关心")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.6),
         [("AI 写代码已经很强了,",        {"size": 44, "bold": True, "color": INK}),
          ("但它会骗你。",                  {"size": 44, "bold": True, "color": RED})])
y = Inches(3.4); cw = Inches(3.7); ch = Inches(2.6); gap = Inches(0.25)
x0 = Inches(0.9)
card(s, x0,                cw[0]*0+y, cw, ch,
     ("现实 1", "写了,但不对",
      ["代码跑得起来,逻辑却错了。",
       "表面上"完成了"。"]))
card(s, x0+cw+gap,         y, cw, ch,
     ("现实 2", "自我感觉良好",
      ["AI 会自信地说:"已完成,并通过所有测试。"",
       "可能一个测试都没跑。"]))
card(s, x0+(cw+gap)*2,     y, cw, ch,
     ("现实 3", "重复犯同一个错",
      ["你纠正一次,它改。",
       "下次新会话,它又犯一模一样的错。"]))
add_text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.6),
         [("解法不是"换更强的模型",而是 —— ", {"size": 18, "color": INK_SOFT}),
          ("给它装上缰绳。",                       {"size": 18, "bold": True, "color": INK})],
         )
slide_num(s, 2)

# ---------- 3. Analogy ----------
s = add_slide(PAPER)
eyebrow(s, Inches(0.9), Inches(0.6), "先讲个比喻")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("Harness 原本是什么?", {"size": 44, "bold": True, "color": INK})])
add_text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(0.9),
         [("Harness 的本意是 ", {"size": 18, "color": INK_SOFT}),
          ("马具 / 缰绳。",       {"size": 18, "bold": True, "color": INK}),
          ("它不让马跑得更快,但它做了三件事 ——",
                                  {"size": 18, "color": INK_SOFT, "space_before": 4})])
y = Inches(3.4); cw = Inches(3.7); ch = Inches(2.4); x0 = Inches(0.9); gap = Inches(0.25)
card(s, x0,                  y, cw, ch,
     ("① 引导方向", "",
      ["缰绳告诉马"去哪儿"。",
       "对应 AI:给它上下文和规则。"]))
card(s, x0+cw+gap,           y, cw, ch,
     ("② 及时拉回", "",
      ["跑偏了立刻反馈。",
       "对应 AI:测试、Lint、类型检查。"]))
card(s, x0+(cw+gap)*2,       y, cw, ch,
     ("③ 解放骑手", "",
      ["骑手不用亲自跑。",
       "对应你:不必每行手写,只管设计。"]))
callout(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.9),
        [("💡 重点:", {"size": 13, "bold": True, "color": CALLOUT_BR}),
         ("没有马,缰绳毫无意义。没有缰绳,马乱跑。AI + Harness = 能用。",
          {"size": 13, "color": INK})])
slide_num(s, 3)

# ---------- 4. Definition ----------
s = add_slide(PAPER)
eyebrow(s, Inches(0.9), Inches(0.6), "正式定义")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("什么是 ", {"size": 40, "bold": True, "color": INK}),
          ("Harness Engineering ?", {"size": 40, "bold": True, "color": ACCENT})])
# Quote block
add_rect(s, Inches(0.9), Inches(2.4), Inches(5.6), Inches(3.6), GRAY_LIGHT, radius=0.08)
add_text(s, Inches(1.2), Inches(2.7), Inches(5.0), Inches(3.0),
         [(""AI Agent 中",         {"size": 26, "bold": True, "color": INK}),
          ("除了模型本身",            {"size": 26, "bold": True, "color": INK}),
          ("以外的一切。"",            {"size": 26, "bold": True, "color": INK}),
          ("— Martin Fowler",        {"size": 13, "color": INK_SOFT, "space_before": 18})])
# Bullets on right
add_text(s, Inches(7.0), Inches(2.4), Inches(5.5), Inches(0.4),
         [("具体包括:", {"size": 14, "bold": True, "color": INK_SOFT})])
items = [
    ("• 给 AI 看的", "上下文", "(已有代码、文档)"),
    ("• AI 能调用的", "工具", "(读文件、跑测试、搜网页)"),
    ("• 约束它的", "规则", "(代码风格、架构边界)"),
    ("• 验证输出的", "检查器", "(测试、Lint、类型)"),
    ("• 多个 AI 之间的", "流程编排", ""),
]
yy = Inches(2.95)
for prefix, bold, suffix in items:
    add_text(s, Inches(7.0), yy, Inches(5.8), Inches(0.55),
             [(prefix, {"size": 15, "color": INK}),
              (bold,    {"size": 15, "bold": True, "color": INK}),
              (suffix,  {"size": 15, "color": INK})])
    yy += Inches(0.55)
slide_num(s, 4)

# ---------- 5. Two goals ----------
s = add_slide(TINT_BLUE)
eyebrow(s, Inches(0.9), Inches(0.6), "Harness 的两个目标")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.6),
         [("一个好缰绳,",                  {"size": 40, "bold": True, "color": INK}),
          ("帮 AI 做到这两件事:",          {"size": 40, "bold": True, "color": INK})])
y = Inches(3.0); cw = Inches(5.6); ch = Inches(3.0); gap = Inches(0.3)
card(s, Inches(0.9),         y, cw, ch,
     ("目标 ①", "第一次就写对",
      ["通过事前引导,让 AI 在下笔前就"知道"这个项目的规矩、边界、约定。",
       "好比新员工入职前先看完员工手册,而不是上来就乱改代码。"]),
     accent_color=ACCENT)
card(s, Inches(0.9)+cw+gap,  y, cw, ch,
     ("目标 ②", "能自己发现并改错",
      ["写完后立刻有反馈:测试失败?Lint 报警?类型不对?AI 自己读报告,自己改。",
       "好比写完作业先自己对答案,再交给老师。"]),
     accent_color=PURPLE)
add_text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.6),
         [("两件事合起来,让 ", {"size": 17, "color": INK_SOFT}),
          ("人 ",                 {"size": 17, "bold": True, "color": INK}),
          ("从"一行行审"解放出来,去做"设计"。",
                                  {"size": 17, "color": INK_SOFT})])
slide_num(s, 5)

# ---------- 6. Two pillars ----------
s = add_slide(PAPER)
eyebrow(s, Inches(0.9), Inches(0.6), "两根支柱")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.6),
         [("Guides & Sensors", {"size": 44, "bold": True, "color": INK}),
          ("引导 + 反馈",        {"size": 22, "color": INK_SOFT, "space_before": 6})])
y = Inches(3.0); cw = Inches(5.6); ch = Inches(2.5); gap = Inches(0.3)
card(s, Inches(0.9),         y, cw, ch,
     ("Feedforward · 前馈", "Guides 引导",
      ["动手之前给 AI 的"地图"。",
       "示例代码 / 代码规范 / README / 架构文档 / 类似场景的过往实现。",
       "→ 让它少犯错。"]))
card(s, Inches(0.9)+cw+gap,  y, cw, ch,
     ("Feedback · 反馈", "Sensors 传感器",
      ["动手之后给 AI 的"体检报告"。",
       "单元测试 / Lint / 类型检查 / 另一个 AI 评审 / 生产监控告警。",
       "→ 让它能改错。"]))
callout(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.4),
        [("⚠️ 只有一根支柱的后果:",
          {"size": 13, "bold": True, "color": CALLOUT_BR, "space_after": 4}),
         ("只有引导 → AI 不知道规则有没有被遵守,重复犯错。",
          {"size": 12, "color": INK, "space_after": 2}),
         ("只有反馈 → AI 每次都从零摸索,效率很低。两个都要,缺一不可。",
          {"size": 12, "color": INK})])
slide_num(s, 6)

# ---------- 7. Guides examples ----------
s = add_slide(GREEN_SOFT)
eyebrow(s, Inches(0.9), Inches(0.6), "Guides 具体长什么样")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.9),
         [(""引导" 的日常例子", {"size": 38, "bold": True, "color": INK})])
y = Inches(2.4); cw = Inches(5.6); ch = Inches(2.2); gap = Inches(0.3)
examples = [
    (Inches(0.9),                y,         "例子 1 · 项目说明文档", "CLAUDE.md / README.md",
     "写清楚:本项目用什么框架、命名约定、禁止哪些依赖、测试怎么跑。AI 每次会话开头自动读一次。"),
    (Inches(0.9)+cw+gap,         y,         "例子 2 · 示例代码",       "类似的旧功能放在它能看到的地方",
     ""按 src/features/login 的风格来实现 signup。" AI 会模仿你已有的好样子。"),
    (Inches(0.9),                y+ch+gap,  "例子 3 · 结构化提示",     "任务模板",
     "先写清楚"需求 / 边界 / 不要做 / 完成标准"再让 AI 动手,比"帮我写个登录"稳定得多。"),
    (Inches(0.9)+cw+gap,         y+ch+gap,  "例子 4 · 工具约束",       "只开放必要工具",
     "不需要网络的任务就别给它联网权限。减少它"自作主张"的空间。"),
]
for x, yy, label, title, body in examples:
    card(s, x, yy, cw, ch, (label, title, [body]))
slide_num(s, 7)

# ---------- 8. Sensors examples ----------
s = add_slide(TINT_WARM)
eyebrow(s, Inches(0.9), Inches(0.6), "Sensors 具体长什么样")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.9),
         [(""反馈" 的日常例子", {"size": 38, "bold": True, "color": INK})])
y = Inches(2.4); cw = Inches(5.6); ch = Inches(2.2); gap = Inches(0.3)
examples = [
    (Inches(0.9),                y,         "例子 1 · 自动测试",     "pytest / jest / go test",
     "AI 写完代码立刻跑测试,看到红色立刻改。最便宜、最可靠的缰绳。"),
    (Inches(0.9)+cw+gap,         y,         "例子 2 · 类型检查",     "TypeScript / mypy / Pyright",
     "编译器报错就是 AI 的第一道"老师"。强类型语言天然适合 harness。"),
    (Inches(0.9),                y+ch+gap,  "例子 3 · Linter",       "ESLint / Ruff / Prettier",
     "自动修复风格问题,AI 不需要去猜"我该用几个空格"。"),
    (Inches(0.9)+cw+gap,         y+ch+gap,  "例子 4 · 真浏览器验证", "Playwright / Puppeteer",
     "对前端,让 AI 真去点页面、截图、看控制台错误。不是"编译过了 = 对的"。"),
]
for x, yy, label, title, body in examples:
    card(s, x, yy, cw, ch, (label, title, [body]))
slide_num(s, 8)

# ---------- 9. Computational vs Inferential ----------
s = add_slide(PAPER)
eyebrow(s, Inches(0.9), Inches(0.6), "反馈分两种")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("确定的 vs. 猜的", {"size": 44, "bold": True, "color": INK})])
y = Inches(2.5); cw = Inches(5.6); ch = Inches(3.4); gap = Inches(0.3)
card(s, Inches(0.9), y, cw, ch,
     ("Computational · 计算式", "✓ 确定的检查",
      ["速度: 毫秒 ~ 秒级",
       "稳定性: 同样输入 = 同样结果",
       "例子: 单元测试、类型检查、Lint、AST 分析、代码覆盖率。",
       "→ 能用这种,优先用这种。"]),
     bg=GREEN_SOFT, accent_color=RGBColor(0x16, 0x65, 0x34))
card(s, Inches(0.9)+cw+gap, y, cw, ch,
     ("Inferential · 推理式", "~ 猜的检查",
      ["速度: 几秒 ~ 几分钟",
       "稳定性: 同样输入可能给出不同结果",
       "例子: 让另一个 AI 评审"这代码是否过度设计了"、"命名好不好"。",
       "⚠️ 贵、慢、不一定准 —— 用在必要的地方。"]),
     bg=RED_SOFT, accent_color=RGBColor(0xb9, 0x1c, 0x1c))
callout(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.9),
        [("初学者口诀:能用确定方法的,绝不让 AI 来猜。",
          {"size": 13, "bold": True, "color": CALLOUT_BR}),
         ("写一个单元测试,比让另一个 AI 来看 10 遍,都更可靠。",
          {"size": 12, "color": INK})])
slide_num(s, 9)

# ---------- 10. Three regulation categories ----------
s = add_slide(PAPER)
eyebrow(s, Inches(0.9), Inches(0.6), "缰绳要管的三件事")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("管什么?", {"size": 44, "bold": True, "color": INK})])
y = Inches(2.5); cw = Inches(3.7); ch = Inches(3.4); gap = Inches(0.25)
x0 = Inches(0.9)
card(s, x0, y, cw, ch,
     ("维度 ①", "可维护性",
      ["内部代码质量。",
       "有没有重复?函数是不是太长?测试够不够?命名清不清楚?",
       "✓ 可以自动化,成熟工具多。"]))
card(s, x0+cw+gap, y, cw, ch,
     ("维度 ②", "架构契合度",
      ["符不符合系统设计。",
       "性能指标、模块边界、安全规范、可观测性。",
       "~ 部分可自动化(Fitness Functions)。"]))
card(s, x0+(cw+gap)*2, y, cw, ch,
     ("维度 ③", "行为正确",
      ["功能到底对不对。",
       "用户登录能登进去吗?付款金额算对了吗?",
       "✕ 最难,目前没有银弹。"]))
add_text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.6),
         [("第三个是最重要也最难的 —— 这是目前 harness engineering 最前沿的研究方向。",
           {"size": 14, "color": INK_SOFT})])
slide_num(s, 10)

# ---------- 11. Steering Loop ----------
s = add_slide(TINT_BLUE)
eyebrow(s, Inches(0.9), Inches(0.6), "人和 AI 怎么协作")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("掌舵循环 ", {"size": 40, "bold": True, "color": INK}),
          ("The Steering Loop", {"size": 40, "bold": True, "color": ACCENT})])
add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.6),
         [("人不写每一行代码,而是", {"size": 18, "color": INK_SOFT}),
          ("持续调整缰绳本身。",     {"size": 18, "bold": True, "color": INK})])
# Diagram of 4 nodes
nodes = [
    ("① 观察 AI 输出", WHITE,  INK),
    ("② 发现重复模式", WHITE,  INK),
    ("③ 改进 Guide / Sensor", WHITE, ACCENT),
    ("④ 下一轮更准", INK, WHITE),
]
nx = Inches(0.9); ny = Inches(3.2); nw = Inches(2.6); nh = Inches(1.4); ngap = Inches(0.25)
for i, (t, fill, txt) in enumerate(nodes):
    px = nx + (nw + ngap) * i
    add_rect(s, px, ny, nw, nh, fill,
             line=ACCENT if txt == ACCENT else (LINE if fill == WHITE else INK), radius=0.10)
    add_text(s, px, ny, nw, nh,
             [(t, {"size": 16, "bold": True, "color": txt})],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        ax = px + nw + Inches(0.04)
        add_text(s, ax, ny, ngap - Inches(0.08), nh,
                 [("→", {"size": 24, "color": INK_SOFT})],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
callout(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.5),
        [("关键心态:", {"size": 13, "bold": True, "color": CALLOUT_BR, "space_after": 3}),
         ("✕ 不要一遍遍手改 AI 的错。", {"size": 12, "color": INK, "space_after": 2}),
         ("✓ 把"你怎么发现它错了"这件事,编码进缰绳,让 AI 下次自己发现。",
          {"size": 12, "color": INK})])
slide_num(s, 11)

# ---------- 12. Keep quality left ----------
s = add_slide(PAPER)
eyebrow(s, Inches(0.9), Inches(0.6), "什么时候检查?")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("越早越便宜", {"size": 44, "bold": True, "color": INK})])
add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
         [("同样一个 bug,发现得越晚,修复代价越大。",
           {"size": 16, "color": INK_SOFT})])
y = Inches(2.9); cw = Inches(2.85); ch = Inches(2.4); gap = Inches(0.2)
steps = [
    ("1", "AI 写代码时", "类型检查、Lint(毫秒)"),
    ("2", "提交前",       "跑相关单元测试(秒)"),
    ("3", "合并前 CI",    "全量测试 + 集成(分钟)"),
    ("4", "上线后",       "生产监控告警(持续)"),
]
for i, (n, t, p) in enumerate(steps):
    px = Inches(0.9) + (cw + gap) * i
    add_rect(s, px, y, cw, ch, WHITE, line=LINE, radius=0.08)
    # Number circle
    add_rect(s, px + Inches(0.3), y + Inches(0.3), Inches(0.5), Inches(0.5), INK, radius=0.5)
    add_text(s, px + Inches(0.3), y + Inches(0.3), Inches(0.5), Inches(0.5),
             [(n, {"size": 14, "bold": True, "color": WHITE})],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, px + Inches(0.3), y + Inches(1.0), cw - Inches(0.6), Inches(1.4),
             [(t, {"size": 16, "bold": True, "color": INK, "space_after": 4}),
              (p, {"size": 12, "color": INK_SOFT, "line_spacing": 1.4})])
callout(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.2),
        [("💡 Keep quality left · 向左看齐:",
          {"size": 13, "bold": True, "color": CALLOUT_BR, "space_after": 3}),
         ("快的检查前置,慢的检查后置。让 AI 在最便宜的阶段自己修掉大部分问题。",
          {"size": 12, "color": INK})])
slide_num(s, 12)

# ---------- 13. Harnessability ----------
s = add_slide(GREEN_SOFT)
eyebrow(s, Inches(0.9), Inches(0.6), "一个冷知识")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.5),
         [("不是所有项目",        {"size": 38, "bold": True, "color": INK}),
          ("都"能装缰绳"",        {"size": 38, "bold": True, "color": INK})])
add_text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.5),
         [("Martin Fowler 称之为 ", {"size": 16, "color": INK_SOFT}),
          ("Harnessability · 可缰绳化程度。",
                                       {"size": 16, "bold": True, "color": INK})])
y = Inches(3.2); cw = Inches(5.6); ch = Inches(2.6); gap = Inches(0.3)
card(s, Inches(0.9), y, cw, ch,
     ("容易装缰绳", "Greenfield · 新项目",
      ["• 强类型语言(TS / Rust / Go)",
       "• 模块边界清晰",
       "• 测试覆盖率高",
       "• 用主流框架"]),
     accent_color=RGBColor(0x16, 0x65, 0x34))
card(s, Inches(0.9)+cw+gap, y, cw, ch,
     ("难装缰绳", "Legacy · 遗留系统",
      ["• 弱类型、动态语言",
       "• 耦合严重,改一处炸一片",
       "• 没测试,或测试跑不起来",
       "• 自研框架,AI 没见过"]),
     accent_color=RGBColor(0xb9, 0x1c, 0x1c))
callout(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(1.0),
        [("⚠️ 悖论:",
          {"size": 13, "bold": True, "color": CALLOUT_BR, "space_after": 3}),
         ("最需要 AI 帮忙的遗留系统,恰恰最难装缰绳。先花时间补测试、补类型,往往是让 AI 好用的前置投资。",
          {"size": 12, "color": INK})])
slide_num(s, 13)

# ---------- 14. Quote 1 ----------
s = add_slide(BLACK)
eyebrow(s, Inches(0.9), Inches(0.6), "Fowler 的一句关键话", dark=True)
add_text(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(4.5),
         [(""Harness 不是替代人类判断,",  {"size": 38, "bold": True, "color": WHITE}),
          ("而是把人类判断",                  {"size": 38, "bold": True, "color": WHITE}),
          ("引导到最关键的地方。"",          {"size": 38, "bold": True, "color": PURPLE})])
add_text(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.5),
         [("— Martin Fowler · Harness Engineering",
           {"size": 14, "color": GRAY_TEXT})])
slide_num(s, 14, dark=True)

# ---------- 15. Part II transition ----------
s = add_slide(BLACK)
eyebrow(s, Inches(0.9), Inches(0.6), "Part II", dark=True)
add_text(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(3.0),
         [("工程师", {"size": 64, "bold": True, "color": WHITE}),
          ("怎么 ", {"size": 64, "bold": True, "color": WHITE, "space_before": 0}),
          ("实际在做?", {"size": 64, "bold": True, "color": PURPLE})])
add_text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.5),
         [("来看 Anthropic 工程师的真实案例 ——",
           {"size": 20, "color": GRAY_TEXT}),
          ("他们怎么让 AI 连续干活 6 个小时不崩。",
           {"size": 20, "color": GRAY_TEXT})])
slide_num(s, 15, dark=True)

# ---------- 16. Long-running problems ----------
s = add_slide(PAPER)
eyebrow(s, Inches(0.9), Inches(0.6), "当任务要跑几小时")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("两个经典"老毛病"", {"size": 42, "bold": True, "color": INK})])
y = Inches(2.6); cw = Inches(5.6); ch = Inches(3.6); gap = Inches(0.3)
card(s, Inches(0.9), y, cw, ch,
     ("问题 ① · 上下文焦虑", "Context Anxiety",
      ["AI 的"工作记忆"有上限。",
       "快满的时候,它会开始草草收尾 —— 明明任务没做完,它会说"做完了,建议后续优化……"",
       "就像人快下班时,会把剩下的活随便糊弄一下。"]))
card(s, Inches(0.9)+cw+gap, y, cw, ch,
     ("问题 ② · 自我表扬", "Sycophancy",
      ["AI 倾向于说自己做得很好,不管实际如何。",
       "让它自己审自己的代码,它会说"优秀 9/10"。",
       "就像让学生自己判自己的卷子。"]))
slide_num(s, 16)

# ---------- 17. Context reset ----------
s = add_slide(TINT_BLUE)
eyebrow(s, Inches(0.9), Inches(0.6), "对付"上下文焦虑"")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("上下文 ", {"size": 40, "bold": True, "color": INK}),
          ("重置",      {"size": 40, "bold": True, "color": ACCENT}),
          (",不是压缩", {"size": 40, "bold": True, "color": INK})])
y = Inches(2.5); cw = Inches(5.6); ch = Inches(3.4); gap = Inches(0.3)
card(s, Inches(0.9), y, cw, ch,
     ("压缩 · compaction", "✕ 让 AI 自己总结历史",
      ["AI 把前面聊的内容压成摘要,继续干。听起来合理,但 ——",
       "摘要会漏掉关键信息,越干越偏,错误在悄悄累积。"]),
     bg=RED_SOFT, accent_color=RGBColor(0xb9, 0x1c, 0x1c))
card(s, Inches(0.9)+cw+gap, y, cw, ch,
     ("重置 · reset", "✓ 结构化交接后开新会话",
      ["写下明确的"已完成 / 下一步 / 重要决定",扔掉旧会话,新 AI 从这份交接文档开始。",
       "像接力赛,每棒都是满血状态。"]),
     bg=GREEN_SOFT, accent_color=RGBColor(0x16, 0x65, 0x34))
callout(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(1.0),
        [("💡 给初学者:",
          {"size": 13, "bold": True, "color": CALLOUT_BR, "space_after": 3}),
         ("你在用 Claude Code / Cursor 聊久了感觉"它变笨了" —— 开个新会话,把关键决定贴进去,通常立刻变聪明。",
          {"size": 12, "color": INK})])
slide_num(s, 17)

# ---------- 18. Generator-Evaluator ----------
s = add_slide(TINT_WARM)
eyebrow(s, Inches(0.9), Inches(0.6), "对付"自我表扬"")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("把 AI ", {"size": 40, "bold": True, "color": INK}),
          ("分成两个角色", {"size": 40, "bold": True, "color": ACCENT})])
# Two big nodes with arrow
nx = Inches(2.5); ny = Inches(2.5); nw = Inches(3.6); nh = Inches(1.6)
add_rect(s, nx, ny, nw, nh, INK, radius=0.10)
add_text(s, nx, ny, nw, nh,
         [("Generator", {"size": 24, "bold": True, "color": WHITE}),
          ("生成器 · 负责写", {"size": 12, "color": GRAY_TEXT})],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, nx + nw, ny, Inches(0.7), nh,
         [("⇄", {"size": 32, "color": INK_SOFT})],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
ex = nx + nw + Inches(0.7)
add_rect(s, ex, ny, nw, nh, WHITE, line=ACCENT, radius=0.10)
add_text(s, ex, ny, nw, nh,
         [("Evaluator", {"size": 24, "bold": True, "color": ACCENT}),
          ("评审员 · 负责挑刺", {"size": 12, "color": INK_SOFT})],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
y = Inches(4.5); cw = Inches(5.6); ch = Inches(2.0); gap = Inches(0.3)
card(s, Inches(0.9), y, cw, ch,
     ("为什么有效", "",
      ["两个 AI 拿不同的 prompt、不同的身份启动,互不串通。",
       "评审员被专门调教得"很挑剔",不会轻易说 OK。"]))
card(s, Inches(0.9)+cw+gap, y, cw, ch,
     ("灵感来源", "",
      ["GAN(生成对抗网络):一个画家,一个鉴赏家。",
       "画家努力骗过鉴赏家,鉴赏家努力挑刺。双方越来越强。"]))
slide_num(s, 18)

# ---------- 19. 4 grading criteria ----------
s = add_slide(PAPER)
eyebrow(s, Inches(0.9), Inches(0.6), "评审员审什么?")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.5),
         [("4 条评分标准",   {"size": 38, "bold": True, "color": INK}),
          ("以"前端设计"任务为例", {"size": 16, "color": INK_SOFT, "space_before": 4})])
y = Inches(2.9); cw = Inches(2.85); ch = Inches(2.6); gap = Inches(0.2)
crits = [
    ("① Design",        "设计品质",   "视觉有没有统一风格?像不像一个团队做的?"),
    ("② Originality",   "原创性",     "是不是又出现了那种"一眼 AI 味"的模板?"),
    ("③ Craft",         "工艺水平",   "对齐、间距、响应式,这些技术基本功扎实吗?"),
    ("④ Functionality", "能不能用",   "按钮真的能点吗?流程真的跑得通吗?"),
]
for i, (label, title, body) in enumerate(crits):
    px = Inches(0.9) + (cw + gap) * i
    card(s, px, y, cw, ch, (label, title, [body]))
add_text(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.0),
         [("评审员用 Playwright 真的去操作页面,不是看着代码猜。",
           {"size": 14, "color": INK_SOFT, "space_after": 4}),
          ("迭代 5 ~ 15 轮,有时会出现"创造性飞跃" —— 换一整套设计方向。",
           {"size": 14, "color": INK_SOFT})])
slide_num(s, 19)

# ---------- 20. 3-agent architecture ----------
s = add_slide(TINT_BLUE)
eyebrow(s, Inches(0.9), Inches(0.6), "更复杂的任务:做整个 App")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("三个 Agent 接力", {"size": 42, "bold": True, "color": INK})])
y = Inches(2.5); cw = Inches(3.7); ch = Inches(3.0); gap = Inches(0.25)
steps = [
    ("1", "Planner · 规划师",  "把"做个音乐制作工具"这种一句话需求,展开成详细规格:有哪些功能、数据模型、界面流程。"),
    ("2", "Generator · 生成器", "按规格逐个 sprint 实现,边写边自查。每次只关注一小块。"),
    ("3", "Evaluator · 评审员", "用 Playwright 真去测功能,写测试报告:"登录通了,支付时金额错了"。"),
]
for i, (n, t, p) in enumerate(steps):
    px = Inches(0.9) + (cw + gap) * i
    add_rect(s, px, y, cw, ch, WHITE, line=LINE, radius=0.10)
    add_rect(s, px + Inches(0.3), y + Inches(0.3), Inches(0.5), Inches(0.5), INK, radius=0.5)
    add_text(s, px + Inches(0.3), y + Inches(0.3), Inches(0.5), Inches(0.5),
             [(n, {"size": 14, "bold": True, "color": WHITE})],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, px + Inches(0.3), y + Inches(1.0), cw - Inches(0.6), Inches(2.0),
             [(t, {"size": 16, "bold": True, "color": INK, "space_after": 6}),
              (p, {"size": 12, "color": INK_SOFT, "line_spacing": 1.45})])
callout(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(1.1),
        [("关键机制 · Sprint Contract:",
          {"size": 13, "bold": True, "color": CALLOUT_BR, "space_after": 3}),
         ("Generator 和 Evaluator 在开干前,先一起把"什么叫做完了"写清楚。",
          {"size": 12, "color": INK})])
slide_num(s, 20)

# ---------- 21. Comparison ----------
s = add_slide(BLACK)
eyebrow(s, Inches(0.9), Inches(0.6), "差距有多大", dark=True)
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("同一道题:做个复古小游戏", {"size": 36, "bold": True, "color": WHITE})])
# Table
rows = [
    ("维度",            "单一 Agent",   "完整 Harness",   None),
    ("耗时",            "20 分钟",       "6 小时",          None),
    ("成本",            "$9",             "$200",            None),
    ("核心玩法",        "✕ 坏的",         "✓ 能玩",          ("red", "green")),
    ("界面打磨",        "✕ 粗糙",         "✓ 精致",          ("red", "green")),
    ("AI 功能集成",     "✕ 没接好",       "✓ 已集成",        ("red", "green")),
    ("用户能直接用吗",  "✕ 不能",         "✓ 能",            ("red", "green")),
]
ty = Inches(2.3); rh = Inches(0.45)
for i, row in enumerate(rows):
    if len(row) == 4:
        col1, col2, col3, marks = row
    else:
        col1, col2, col3 = row; marks = None
    yy = ty + rh * i
    if i == 0:
        # Header row
        add_text(s, Inches(0.9),  yy, Inches(4.0), rh,
                 [(col1, {"size": 12, "bold": True, "color": GRAY_TEXT})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(5.0),  yy, Inches(3.5), rh,
                 [(col2, {"size": 12, "bold": True, "color": GRAY_TEXT})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.7),  yy, Inches(3.5), rh,
                 [(col3, {"size": 12, "bold": True, "color": GRAY_TEXT})],
                 anchor=MSO_ANCHOR.MIDDLE)
    else:
        c2_color = RED if marks and marks[0] == "red" else WHITE
        c3_color = GREEN if marks and marks[1] == "green" else WHITE
        add_text(s, Inches(0.9),  yy, Inches(4.0), rh,
                 [(col1, {"size": 14, "color": WHITE})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(5.0),  yy, Inches(3.5), rh,
                 [(col2, {"size": 14, "color": c2_color})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.7),  yy, Inches(3.5), rh,
                 [(col3, {"size": 14, "color": c3_color})],
                 anchor=MSO_ANCHOR.MIDDLE)
    add_line(s, Inches(0.9), yy + rh, Inches(12.4), yy + rh, color=DARK_CARD)
add_text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.7),
         [("花多 22 倍的钱,不是浪费 —— 是把"能跑的 demo"变成"真能用的产品"。",
           {"size": 14, "color": GRAY_TEXT})])
slide_num(s, 21, dark=True)

# ---------- 22. Evolution ----------
s = add_slide(PAPER)
eyebrow(s, Inches(0.9), Inches(0.6), "有趣的反转")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.5),
         [("模型变强后,",        {"size": 40, "bold": True, "color": INK}),
          ("缰绳要",                {"size": 40, "bold": True, "color": INK}),
          ("变松", {"size": 40, "bold": True, "color": ACCENT})])
add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.6),
         [("Opus 4.6 / 4.7 出来后,Anthropic 的工程师主动拆掉了一些脚手架:",
           {"size": 17, "color": INK_SOFT})])
items = [
    "• 以前要切成多个小 sprint → 现在一次性做 2 小时没问题。",
    "• 以前 Evaluator 要反复挑刺 → 现在"一次过审"就够了。",
    "• 以前要三个 agent → 现在很多任务一个就搞定。",
]
yy = Inches(3.5)
for it in items:
    add_text(s, Inches(0.9), yy, Inches(11.5), Inches(0.5),
             [(it, {"size": 17, "color": INK})])
    yy += Inches(0.55)
callout(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.5),
        [("💡 原则:",
          {"size": 13, "bold": True, "color": CALLOUT_BR, "space_after": 3}),
         ("用最简单能跑通的 harness。Harness 是脚手架,不是建筑。",
          {"size": 13, "color": INK, "space_after": 3}),
         ("不要为未来可能的问题过度设计 —— 那是工程师的通病,在 AI 时代代价更大。",
          {"size": 12, "color": INK})])
slide_num(s, 22)

# ---------- 23. Ashby's Law ----------
s = add_slide(TINT_WARM)
eyebrow(s, Inches(0.9), Inches(0.6), "一个深层原理")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.5),
         [("Ashby's Law", {"size": 44, "bold": True, "color": INK}),
          ("必要变量定律",  {"size": 22, "color": INK_SOFT, "space_before": 6})])
add_text(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(0.9),
         [("控制论里的一句老话:",
           {"size": 16, "color": INK_SOFT, "space_after": 4}),
          (""控制器的复杂度必须能匹配被控系统的复杂度。"",
           {"size": 18, "bold": True, "color": INK})])
y = Inches(3.9); cw = Inches(5.6); ch = Inches(2.0); gap = Inches(0.3)
card(s, Inches(0.9), y, cw, ch,
     (None, "翻译成人话",
      ["AI 能写出的花样越多,你的缰绳就要越复杂,才能管住。"]))
card(s, Inches(0.9)+cw+gap, y, cw, ch,
     (None, "聪明的反向操作",
      ["先主动减少花样 —— 规定好"就用这三种架构模板之一",AI 的输出空间变小,缰绳就容易做。"]),
     accent_color=ACCENT)
callout(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(1.0),
        [("初学者类比:",
          {"size": 13, "bold": True, "color": CALLOUT_BR, "space_after": 3}),
         ("教小孩画画,给他一盒 8 色蜡笔比给他 120 色更好教 —— 选择少了,你好管、他好学。",
          {"size": 12, "color": INK})])
slide_num(s, 23)

# ---------- 24. Harness templates ----------
s = add_slide(PAPER)
eyebrow(s, Inches(0.9), Inches(0.6), "企业可能的下一步")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("Harness 模板化", {"size": 42, "bold": True, "color": INK})])
add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.6),
         [("不同类型的服务,缰绳长得不一样。未来可能会出现 ", {"size": 16, "color": INK_SOFT}),
          ("预配置的 harness 包",
           {"size": 16, "bold": True, "color": INK})])
y = Inches(3.0); cw = Inches(3.7); ch = Inches(2.4); gap = Inches(0.25)
templates = [
    ("模板 A", "CRUD 业务服务", "一套固定的:数据模型 / API 约定 / 测试模板 / 性能阈值。"),
    ("模板 B", "事件处理器",     "重点在幂等性、消息格式校验、重试机制的 harness。"),
    ("模板 C", "数据看板",       "重点在查询性能、图表规范、数据新鲜度检查。"),
]
for i, (label, title, body) in enumerate(templates):
    px = Inches(0.9) + (cw + gap) * i
    card(s, px, y, cw, ch, (label, title, [body]))
add_text(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.9),
         [("未来团队选技术栈,可能会优先考虑 "这个技术有没有现成 harness 模板"",
           {"size": 14, "color": INK_SOFT, "space_after": 3}),
          ("—— 而不只是"这个技术本身好不好"。",
           {"size": 14, "color": INK_SOFT})])
slide_num(s, 24)

# ---------- 25. Open questions ----------
s = add_slide(TINT_BLUE)
eyebrow(s, Inches(0.9), Inches(0.6), "诚实地说")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("还没解决的难题", {"size": 42, "bold": True, "color": INK})])
items = [
    "• 缰绳越装越多,彼此会不会打架?怎么保证一致?",
    "• 怎么衡量一个 harness 好不好?覆盖度怎么算?",
    "• 反馈太多,AI 会被噪音淹没 —— 怎么过滤?",
    "• 最大的:怎么真正验证"行为正确"?测试只能覆盖想得到的场景。",
]
yy = Inches(2.4)
for it in items:
    add_text(s, Inches(0.9), yy, Inches(11.5), Inches(0.6),
             [(it, {"size": 18, "color": INK})])
    yy += Inches(0.7)
callout(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.4),
        [("这些问题目前没有标准答案。",
          {"size": 13, "bold": True, "color": CALLOUT_BR, "space_after": 3}),
         ("但这也是机会 —— Harness Engineering 是 2026 年最新、最热的软件工程方向之一。",
          {"size": 12, "color": INK})])
slide_num(s, 25)

# ---------- 26. Role of humans ----------
s = add_slide(PAPER)
eyebrow(s, Inches(0.9), Inches(0.6), "你的角色变了")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("人还需要做什么?", {"size": 42, "bold": True, "color": INK})])
y = Inches(2.5); cw = Inches(5.6); ch = Inches(4.2); gap = Inches(0.3)
card(s, Inches(0.9), y, cw, ch,
     ("AI 做不到的", "人类的"隐性 Harness"",
      ["判断力:这个需求的背后,客户真正想要的是什么?",
       "上下文:老板上周才说的政策变化、团队刚踩的坑、法务的边界。",
       "品味:这个 UI 是不是"有灵魂" —— AI 还不太擅长。"]))
card(s, Inches(0.9)+cw+gap, y, cw, ch,
     ("你的新工作", "Harness 设计师",
      ["把你的判断变成规则、模板、检查器。",
       "把踩过的坑写进 CLAUDE.md。",
       "把团队的品味固化成示例代码。",
       "→ 你不是被 AI 替代,你是它的驯兽师。"]),
     accent_color=ACCENT)
slide_num(s, 26)

# ---------- 27. Beginner checklist ----------
s = add_slide(GREEN_SOFT)
eyebrow(s, Inches(0.9), Inches(0.6), "今天就能做的 5 件事")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("新手起步清单", {"size": 42, "bold": True, "color": INK})])
items = [
    "在项目根目录写一个 CLAUDE.md / AGENTS.md,放项目约定。AI 会自动读。",
    "把 自动测试 跑起来 —— 哪怕只有一个。AI 写完代码后让它自己跑。",
    "打开 类型检查(TS / Pyright / mypy)。这是最便宜的缰绳。",
    "碰到 AI 重复犯同一个错,不要手动改第二次 —— 回去改 Guide。",
    "长任务做到一半感觉 AI "变笨"了:开新会话,把关键决定贴过去。",
]
yy = Inches(2.5)
for it in items:
    add_text(s, Inches(0.9), yy, Inches(0.6), Inches(0.6),
             [("✓", {"size": 22, "bold": True, "color": GREEN})],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.5), yy, Inches(11.0), Inches(0.6),
             [(it, {"size": 16, "color": INK})],
             anchor=MSO_ANCHOR.MIDDLE)
    add_line(s, Inches(0.9), yy + Inches(0.7), Inches(12.4), yy + Inches(0.7))
    yy += Inches(0.8)
slide_num(s, 27)

# ---------- 28. Anti-patterns ----------
s = add_slide(PAPER)
eyebrow(s, Inches(0.9), Inches(0.6), "新手常踩的坑")
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
         [("这些做法 ", {"size": 42, "bold": True, "color": INK}),
          ("别做",        {"size": 42, "bold": True, "color": RED})])
items = [
    "让 AI 自己审自己的代码。它会说"完美,9/10" —— 没用。",
    "一个会话一直聊到天荒地老。上下文爆了,它只会糊弄。",
    "为每一个可能的未来需求写规则。Harness 不是越多越好,越臃肿越难维护。",
    "认为 AI 写出来能编译就是对的。编译过 ≠ 逻辑对 ≠ 用户能用。",
    "遇到不确定就问 AI。它的回答看起来很自信,但可能是编的。关键决定留给人。",
]
yy = Inches(2.5)
for it in items:
    add_text(s, Inches(0.9), yy, Inches(0.6), Inches(0.6),
             [("✕", {"size": 22, "bold": True, "color": RED})],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.5), yy, Inches(11.0), Inches(0.6),
             [(it, {"size": 16, "color": INK})],
             anchor=MSO_ANCHOR.MIDDLE)
    add_line(s, Inches(0.9), yy + Inches(0.7), Inches(12.4), yy + Inches(0.7))
    yy += Inches(0.8)
slide_num(s, 28)

# ---------- 29. Final quote ----------
s = add_slide(BLACK)
eyebrow(s, Inches(0.9), Inches(0.6), "记住这句", dark=True)
add_text(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(4.0),
         [(""找到",                          {"size": 50, "bold": True, "color": WHITE}),
          ("能跑通的",                          {"size": 50, "bold": True, "color": WHITE}),
          ("最简单的 harness。"",              {"size": 50, "bold": True, "color": PURPLE})])
add_text(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.5),
         [("— Prithvi Rajasekaran · Anthropic",
           {"size": 14, "color": GRAY_TEXT})])
add_text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
         [("不是越复杂越好。够用就好。模型变强,就把脚手架拆掉。",
           {"size": 14, "color": GRAY_TEXT})])
slide_num(s, 29, dark=True)

# ---------- 30. Closing ----------
s = add_slide(BLACK)
eyebrow(s, Inches(0.9), Inches(0.6), "总结", dark=True)
add_text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(3.5),
         [("你的工作",          {"size": 56, "bold": True, "color": WHITE}),
          ("不再是写代码,",       {"size": 56, "bold": True, "color": GRAY_TEXT}),
          ("是设计这套缰绳。",    {"size": 56, "bold": True, "color": PURPLE})])
add_text(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.0),
         [("Harness Engineering 是 2026 年软件工程最值得学的方向之一。",
           {"size": 18, "color": GRAY_TEXT, "space_after": 4}),
          ("从今天开始,给你的 AI 装上第一根缰绳吧。",
           {"size": 18, "color": GRAY_TEXT})])
add_text(s, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.4),
         [("Sources: Martin Fowler — Harness Engineering · Anthropic — Harness Design for Long-Running Apps",
           {"size": 10, "color": GRAY_TEXT})])
slide_num(s, 30, dark=True)


# ---------- Save ----------
out = "/Users/wenfengzhu/ClaudeProjects/Harness/harness-engineering.pptx"
prs.save(out)
print(f"OK: wrote {out} with {len(prs.slides)} slides")
