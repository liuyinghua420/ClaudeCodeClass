#!/usr/bin/env python3
"""Generate Claude Code course outline PPT from design doc."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_MID = RGBColor(0x16, 0x21, 0x3e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
SUBTITLE_GRAY = RGBColor(0xA0, 0xA0, 0xB0)
DIM_GRAY = RGBColor(0x80, 0x80, 0x90)
ACCENT_BLUE = RGBColor(0x3A, 0x86, 0xFF)
GREEN = RGBColor(0x2E, 0xCC, 0x71)


def add_bg(slide):
    """Add dark background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, left, top, width, height, items, font_size=16,
                     color=LIGHT_GRAY, bullet_color=ORANGE, spacing=Pt(8)):
    """Add bulleted text to slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        # Bullet marker
        run_bullet = p.add_run()
        run_bullet.text = "▸ "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Microsoft YaHei"

        # Text content
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Microsoft YaHei"

    return tf


def add_session_slide(prs, session_num, title, subtitle, hour1_items, hour2_items, milestone=None):
    """Template for session detail slides."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)

    # Session number badge
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(1.5), Inches(0.5),
                f"第 {session_num} 次课", font_size=14, color=ORANGE, bold=True)

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
                title, font_size=30, color=WHITE, bold=True)

    # Subtitle
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                subtitle, font_size=16, color=SUBTITLE_GRAY)

    # Hour 1
    add_textbox(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(0.5),
                "第一小时", font_size=18, color=ACCENT_BLUE, bold=True)
    add_bullet_slide(slide, Inches(0.8), Inches(2.9), Inches(5.5), Inches(3.5),
                     hour1_items, font_size=15)

    # Hour 2
    add_textbox(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(0.5),
                "第二小时", font_size=18, color=ACCENT_BLUE, bold=True)
    add_bullet_slide(slide, Inches(7), Inches(2.9), Inches(5.5), Inches(3.5),
                     hour2_items, font_size=15)

    # Milestone
    if milestone:
        add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
                    f"🎯 里程碑：{milestone}", font_size=16, color=GREEN, bold=True)

    return slide


# ============================================================
# Slide 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "Claude Code 实战课程", font_size=44, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
            "面向零基础成年人 · 8 次课 × 2 小时 · 从入门到做出自己的项目",
            font_size=20, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
            "课程大纲 · 教学参考",
            font_size=16, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
            "2026", font_size=14, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: Course Philosophy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "课程理念", font_size=30, color=WHITE, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.8),
            '这不是一门"学编程"的课。这是一门"学会指挥 AI 智能体"的课。',
            font_size=22, color=ORANGE, bold=True)

items = [
    "核心认知转变：你不是在学编程语言，你是在学管理一个极快的初级开发者",
    "它能阅读整个代码库、做决策、执行任务 —— 你是经理，它是执行者",
    "核心技能：描述需求、迭代优化、用 AI 调试、评估输出结果",
    "市场空白：所有 vibe coding 课都在教 写提示词->得到代码",
    '没有人在教智能体管理思维（上下文管理、迭代优化、评估判断）',
    "这才是真正有价值的技能",
]
add_bullet_slide(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(4), items, font_size=17)


# ============================================================
# Slide 3: Course Overview / Roadmap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "课程路线图：8 次课总览", font_size=30, color=WHITE, bold=True)

sessions = [
    ("1", '"哇"时刻 + 环境搭建', "看到 AI 10 分钟做出网页，安装开发环境"),
    ("2", "AI 思维模型 + 初次动手", "三大思维模型，亲手用 Claude Code 做个人页面"),
    ("3", "小项目 1 · 上半场", "引导式构建日程规划器：拆需求、搭框架"),
    ("4", "小项目 1 · 下半场", "添加功能、调试打磨，完成第一个可用工具"),
    ("5", "小项目 2 · 上半场", "从 4 个模板中选择，写简报，半引导式构建"),
    ("6", "小项目 2 · 下半场", "进阶技巧（Plan Mode），完成并展示"),
    ("7", "个人项目 · 你的想法", "创意工坊 + 构建冲刺"),
    ("8", "完成 + 演示日", "打磨项目，2 分钟展示，庆祝成果"),
]

y_start = 1.5
for i, (num, title, desc) in enumerate(sessions):
    y = y_start + i * 0.7
    # Number
    add_textbox(slide, Inches(0.8), Inches(y), Inches(0.5), Inches(0.5),
                num, font_size=20, color=ORANGE, bold=True)
    # Title
    add_textbox(slide, Inches(1.4), Inches(y), Inches(3.5), Inches(0.5),
                title, font_size=17, color=WHITE, bold=True)
    # Description
    add_textbox(slide, Inches(5.2), Inches(y), Inches(7.5), Inches(0.5),
                desc, font_size=15, color=SUBTITLE_GRAY)


# ============================================================
# Slide 4: Teaching Method
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "教学方法：三级脚手架", font_size=30, color=WHITE, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.5),
            '从"手把手带"到"放手让你飞"的渐进式自主能力培养',
            font_size=17, color=SUBTITLE_GRAY)

# Level 1
add_textbox(slide, Inches(0.8), Inches(2.3), Inches(3.5), Inches(0.5),
            "引导式（第 3-4 次课）", font_size=20, color=ORANGE, bold=True)
items1 = [
    "老师逐步带领，先演示后跟做",
    "每 10-15 分钟一个检查点",
    "所有人做同一类型项目",
    "个性化细节保持趣味",
]
add_bullet_slide(slide, Inches(0.8), Inches(3.0), Inches(3.5), Inches(3), items1, font_size=14)

# Level 2
add_textbox(slide, Inches(4.8), Inches(2.3), Inches(3.5), Inches(0.5),
            "半引导式（第 5-6 次课）", font_size=20, color=ORANGE, bold=True)
items2 = [
    "学员从模板出发独立完成",
    "模板含说明 + 建议构建顺序",
    "3-4 个里程碑检查点",
    "老师巡场，按需帮助",
]
add_bullet_slide(slide, Inches(4.8), Inches(3.0), Inches(3.5), Inches(3), items2, font_size=14)

# Level 3
add_textbox(slide, Inches(8.8), Inches(2.3), Inches(3.8), Inches(0.5),
            "独立式（第 7-8 次课）", font_size=20, color=ORANGE, bold=True)
items3 = [
    "学员构思并构建自己的项目",
    "提供创意菜单消除空白页",
    "老师按需辅导",
    "演示日展示成果",
]
add_bullet_slide(slide, Inches(8.8), Inches(3.0), Inches(3.8), Inches(3), items3, font_size=14)


# ============================================================
# Slides 5-12: Session Details
# ============================================================

add_session_slide(prs, 1,
    '"哇"时刻 + 环境搭建',
    "让学员亲眼看到 AI 的威力，然后装好工具准备上手",
    [
        "老师现场演示：征集学员想法，10 分钟做出网页",
        "想法引导：选单页、静态、视觉化的项目",
        "展示完整循环：描述 → AI 构建 → 浏览器看结果",
        "备用方案：2-3 个预测试 demo + 预录屏幕录像",
    ],
    [
        "双语安装指南（打印 + 投影）",
        "安装清单：终端 → Node.js → Claude Code → API Key",
        "2 名助教走动帮助 + USB 离线安装备份",
        "提前完成的学员：试用速查卡命令",
        "备用：无法安装的学员用网页版",
    ])

add_session_slide(prs, 2,
    "AI 思维模型 + 初次动手",
    "建立正确的 AI 认知框架，然后立刻用起来",
    [
        "前 20 分钟：安装完成缓冲时间",
        "思维模型 1：Claude Code = 初级开发者，你 = 经理",
        "思维模型 2：上下文决定一切（演示好坏对比）",
        "思维模型 3：迭代，而非追求完美",
    ],
    [
        "核心循环练习：描述 → 生成 → 评估 → 优化",
        "练习 1：做一个个人介绍页面",
        "练习 2：修改它（换颜色、加区块、修问题）",
        '核心技能：判断"对"还是"不对，改这里"',
    ])

add_session_slide(prs, 3,
    "小项目 1 —— 引导式构建 · 上半场",
    "所有人做日程规划器，学会拆需求和搭框架",
    [
        "需求思维：把大想法拆成小的可构建部分",
        "练习：用日常语言写 5 条需求描述",
        "喂给 Claude Code，观察结果，评估",
    ],
    [
        "按步骤跟做引导式构建",
        "老师走动帮助解决卡点",
        "关键教学：Claude 做了意外的事怎么纠正",
        "个性化：饮食/运动/工作任务各取所需",
    ])

add_session_slide(prs, 4,
    "小项目 1 —— 引导式构建 · 下半场",
    "添加功能、调试问题、打磨细节",
    [
        "给工具添加 2-3 个新功能",
        "教学：如何给 AI 提供现有代码上下文",
        '实践：多轮对话，"先读这个文件，然后添加..."',
    ],
    [
        "AI 辅助调试：描述期望 vs 实际",
        "学员修复项目中的问题",
        "打磨：让它更好看，添加收尾细节",
    ],
    milestone="每个人都有了一个自己亲手做的可用工具")

add_session_slide(prs, 5,
    "小项目 2 —— 半引导式构建 · 上半场",
    "从 4 个模板中选择，写简报，开始独立构建",
    [
        "模板 A：小型商业着陆页",
        "模板 B：个人作品集/兴趣展示",
        "模板 C：简单互动游戏",
        "模板 D：实用工具（转换器/计时器等）",
        "选择模板，定制需求，写项目简报",
    ],
    [
        "以模板为起点独立构建",
        "老师提供检查点时间线",
        "快慢配对：进度慢的与快的结对",
        "老师审阅简报并提改进建议",
    ])

add_session_slide(prs, 6,
    "小项目 2 —— 半引导式构建 · 下半场",
    "学习进阶技巧，完成项目并展示",
    [
        "进阶：Plan Mode（先规划再构建）",
        "进阶：文件组织和上下文管理",
        "学员将进阶技巧应用到项目中",
    ],
    [
        "完成小项目 2",
        "快速展示：向邻座展示（每人 2-3 分钟）",
        "布置作业：带个人项目想法来第 7 次课",
        "发放项目创意菜单",
    ],
    milestone="学员做出了自己选择的东西，而不仅是跟着做")

add_session_slide(prs, 7,
    "个人项目 —— 你的想法",
    "从想法到原型，老师全程辅导",
    [
        '写下："我的生活/工作中什么能用软件变简单"',
        "提供 10 个预设项目创意菜单",
        "老师帮助精炼范围（5 条写不完就太大）",
        "写项目简报",
    ],
    [
        "构建冲刺：学员构建个人项目",
        "老师逐个辅导，在座位间走动",
        "原则：帮解决卡点，不代做",
        "让学员有效地挣扎",
    ])

add_session_slide(prs, 8,
    "完成 + 演示日",
    "打磨项目，展示成果，庆祝毕业",
    [
        "最后的项目构建和打磨时间",
        "准备 2 分钟演示",
        "（做了什么、怎么做的、学到了什么）",
        "不要求部署，本地演示即可",
    ],
    [
        "每位学员展示项目（每人 2 分钟）",
        "庆祝大家的成果",
        '提供"接下来怎么走"资源',
        "建立班级微信群持续学习",
    ],
    milestone="课程完成！每位学员都有自己的作品")


# ============================================================
# Slide 13: Student Requirements
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "学员要求", font_size=30, color=WHITE, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5),
            "前提条件", font_size=20, color=ORANGE, bold=True)
items = [
    "macOS 10.15+ 或 Windows 10+ 笔记本电脑",
    "拥有管理员/安装权限",
    "稳定网络连接",
    "不需要任何编程经验",
]
add_bullet_slide(slide, Inches(0.8), Inches(2.1), Inches(5), Inches(3), items, font_size=16)

add_textbox(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.5),
            "课程配置", font_size=20, color=ORANGE, bold=True)
items2 = [
    "班级规模：10-15 人",
    "1 名老师 + 2 名助教",
    "中文授课，工具界面英文",
    "讲义双语（中文讲解 + 英文命令）",
    "每节课约 55 分钟休息 10 分钟",
    "技术范围：纯前端 HTML/CSS/JS",
]
add_bullet_slide(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(3.5), items2, font_size=16)


# ============================================================
# Slide 14: Success Criteria
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "成功标准", font_size=30, color=WHITE, bold=True)

metrics = [
    ("95%+", "第 1 次课结束", "学员能执行 Claude Code 命令"),
    ("90%+", "第 4 次课结束", "学员有可工作的小项目 1"),
    ("80%+", "第 8 次课", "学员有可演示的个人项目"),
    ("主观", "课后反馈", '学员表示"我可以自己做了"'),
    ("3+", "课后追踪", "学员继续独立使用 Claude Code"),
]

for i, (num, timing, desc) in enumerate(metrics):
    y = 1.6 + i * 1.0
    add_textbox(slide, Inches(0.8), Inches(y), Inches(1.5), Inches(0.6),
                num, font_size=32, color=ORANGE, bold=True)
    add_textbox(slide, Inches(2.5), Inches(y), Inches(3), Inches(0.5),
                timing, font_size=16, color=SUBTITLE_GRAY)
    add_textbox(slide, Inches(2.5), Inches(y + 0.35), Inches(9), Inches(0.5),
                desc, font_size=18, color=WHITE)


# ============================================================
# Slide 15: Materials Checklist
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "准备材料清单", font_size=30, color=WHITE, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
            "教学材料", font_size=20, color=ORANGE, bold=True)
items = [
    "打印安装指南（覆膜，可重复使用）",
    "速查卡：20 个常用 Claude Code 命令",
    "项目简报模板（填空格式）",
    "故障排除指南",
    "课后资源清单",
    "项目创意菜单（10 个预设想法）",
]
add_bullet_slide(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4), items, font_size=15)

add_textbox(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.5),
            "技术准备", font_size=20, color=ORANGE, bold=True)
items2 = [
    "2-3 个 Anthropic API Key（分组使用）",
    "USB 驱动器：Node.js 安装包 + 离线包",
    "预录演示视频（API 故障备用）",
    "预缓存的 demo 输出",
    "4 个小项目 2 模板文件夹",
    "课前测试 API 并发上限",
]
add_bullet_slide(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(4), items2, font_size=15)


# ============================================================
# Slide 16: Risk Mitigation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "风险与应对", font_size=30, color=WHITE, bold=True)

risks = [
    ("安装失败", "网页版备用 + USB 离线安装 + 第 2 次课缓冲时间"),
    ("API 故障/限速", "预缓存 demo 输出 + 离线练习 + 备用 Key"),
    ("技能差距大", "快慢配对 + 差异化检查点 + 多级难度"),
    ("项目范围膨胀", "5 条规则 + 创意菜单 + 老师否决权"),
    ("学员缺课", "微信发课程回顾 + 追赶检查点 + 预录视频"),
    ("费用顾虑", "老师提供共享 Key，预算约 $300"),
]

for i, (risk, mitigation) in enumerate(risks):
    y = 1.5 + i * 0.9
    add_textbox(slide, Inches(0.8), Inches(y), Inches(2.5), Inches(0.5),
                risk, font_size=17, color=ORANGE, bold=True)
    add_textbox(slide, Inches(3.5), Inches(y), Inches(9), Inches(0.5),
                mitigation, font_size=16, color=LIGHT_GRAY)


# ============================================================
# Slide 17: Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "下一步行动", font_size=30, color=WHITE, bold=True)

steps = [
    ("立即", "创建第 1 次课安装指南，找 2-3 个零基础朋友测试"),
    ("第 1 周", '制作"哇"时刻演示脚本，练到流畅'),
    ("第 2 周", "端到端设计小项目 1（引导式构建），包括所有检查点"),
    ("第 3 周", "创建小项目 2 的 4 个项目模板"),
    ("开课前", "找 3-5 人试讲第 1-2 次课"),
    ("开课前", "验证网页版备用方案 + 准备共享备用电脑"),
]

for i, (timing, action) in enumerate(steps):
    y = 1.5 + i * 0.85
    add_textbox(slide, Inches(0.8), Inches(y), Inches(1.5), Inches(0.5),
                timing, font_size=16, color=ORANGE, bold=True)
    add_textbox(slide, Inches(2.5), Inches(y), Inches(10), Inches(0.5),
                action, font_size=17, color=LIGHT_GRAY)


# ============================================================
# Save
# ============================================================
output_path = "/Users/wenfengzhu/ClaudeProjects/ClaudeClass/ClaudeCode课程大纲.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
