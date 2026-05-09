#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate Lesson 1 teaching PPT from lesson1-spec.md."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
SUBTITLE_GRAY = RGBColor(0xA0, 0xA0, 0xB0)
DIM_GRAY = RGBColor(0x80, 0x80, 0x90)
ACCENT_BLUE = RGBColor(0x3A, 0x86, 0xFF)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
CODE_BG = RGBColor(0x0D, 0x11, 0x17)
YELLOW = RGBColor(0xF1, 0xC4, 0x0F)
RED_SOFT = RGBColor(0xE7, 0x4C, 0x3C)


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_text(slide, left, top, width, height, text, size=18,
             color=LIGHT_GRAY, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.alignment = align
    return tf


def add_bullets(slide, left, top, width, height, items, size=16,
                color=LIGHT_GRAY, bullet_color=ORANGE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r1 = p.add_run()
        r1.text = '\u25b8 '
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.name = 'Microsoft YaHei'
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = 'Microsoft YaHei'
    return tf


def add_code(slide, left, top, width, height, code, size=14):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(size)
    p.font.color.rgb = GREEN
    p.font.name = 'Courier New'
    p.alignment = PP_ALIGN.LEFT
    return tf


def section_slide(title, subtitle=None):
    """Big centered section title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
             title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8),
                 subtitle, size=20, color=SUBTITLE_GRAY, align=PP_ALIGN.CENTER)
    return slide


def teaching_slide(time_tag, title, items_left, items_right=None,
                   note=None, code=None):
    """Standard two-column teaching slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    # Time badge
    add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
             time_tag, size=13, color=ORANGE, bold=True)
    # Title
    add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             title, size=28, color=WHITE, bold=True)

    if items_right:
        add_bullets(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.5),
                    items_left, size=16)
        add_bullets(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(3.5),
                    items_right, size=16)
    else:
        add_bullets(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(3.5),
                    items_left, size=17)

    if code:
        add_code(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(1.5), code)

    if note:
        add_text(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
                 note, size=14, color=YELLOW, bold=True)
    return slide


# ============================================================
# Slide 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(0.6),
         'Claude Code \u5b9e\u6218\u8bfe\u7a0b', size=20, color=SUBTITLE_GRAY,
         align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
         '\u7b2c 1 \u6b21\u8bfe', size=54, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
         '\u201c\u54c7\u201d\u65f6\u523b + \u73af\u5883\u642d\u5efa',
         size=28, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
         '\u4eca\u5929\u7684\u76ee\u6807\uff1a\u4eb2\u773c\u89c1\u8bc1 AI \u7684\u5a01\u529b\uff0c\u88c5\u597d\u5de5\u5177\uff0c\u8dd1\u901a\u7b2c\u4e00\u4e2a\u9879\u76ee',
         size=16, color=DIM_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: Today's Goals
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         '\u4eca\u5929\u7ed3\u675f\u65f6\uff0c\u4f60\u4f1a\u2026\u2026',
         size=30, color=WHITE, bold=True)

goals = [
    ('\u2764\ufe0f', '\u5185\u5fc3\u88ab\u70b9\u71c3',
     '\u4eb2\u773c\u770b\u5230 AI \u5728 10 \u5206\u949f\u5185\u628a\u4e00\u4e2a\u60f3\u6cd5\u53d8\u6210\u53ef\u7528\u7f51\u9875'),
    ('\u2699\ufe0f', '\u73af\u5883\u5c31\u7eea',
     '\u7535\u8111\u4e0a\u80fd\u8dd1\u901a Claude Code \u547d\u4ee4'),
    ('\u2328\ufe0f', '\u5b8c\u6210\u9996\u6b21\u4ea4\u4e92',
     '\u81ea\u5df1\u8ddf Claude Code \u5bf9\u8bdd\u8fc7\u4e00\u6b21\uff0c\u770b\u5230\u8fc7\u5b83\u7684\u8f93\u51fa'),
]
for i, (emoji, title, desc) in enumerate(goals):
    y = 1.8 + i * 1.6
    add_text(slide, Inches(0.8), Inches(y), Inches(0.8), Inches(0.6),
             emoji, size=32, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.8), Inches(y), Inches(4), Inches(0.5),
             title, size=24, color=ORANGE, bold=True)
    add_text(slide, Inches(1.8), Inches(y + 0.5), Inches(10), Inches(0.5),
             desc, size=17, color=LIGHT_GRAY)


# ============================================================
# Slide 3: Timeline Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         '\u4eca\u5929\u7684\u65f6\u95f4\u7ebf', size=30, color=WHITE, bold=True)

timeline = [
    ('0:00', '5 min', '\u5f00\u573a', '\u81ea\u6211\u4ecb\u7ecd\uff0c\u8bfe\u7a0b\u6982\u8ff0'),
    ('0:05', '5 min', '\u94fa\u57ab', '\u4ec0\u4e48\u662f Claude Code\uff1f'),
    ('0:10', '10 min', '\u201c\u54c7\u201d\u65f6\u523b\u6f14\u793a',
     '\u73b0\u573a\u5f81\u96c6\u60f3\u6cd5\uff0c10 \u5206\u949f\u505a\u51fa\u7f51\u9875'),
    ('0:20', '10 min', '\u62c6\u89e3\u6f14\u793a',
     '\u89e3\u91ca\u6838\u5fc3\u5faa\u73af\uff1a\u63cf\u8ff0 \u2192 \u751f\u6210 \u2192 \u8bc4\u4f30 \u2192 \u4f18\u5316'),
    ('0:30', '10 min', '\u8bfe\u7a0b\u8def\u7ebf\u56fe',
     '8 \u6b21\u8bfe\u7684\u65c5\u7a0b\u603b\u89c8'),
    ('0:40', '15 min', '\u7ec8\u7aef\u5165\u95e8',
     '\u6253\u5f00\u7ec8\u7aef\uff0c\u5b66 5 \u4e2a\u57fa\u7840\u547d\u4ee4'),
    ('0:55', '10 min', '\u4f11\u606f', '\u2615'),
    ('1:05', '25 min', '\u73af\u5883\u5b89\u88c5',
     'Node.js \u2192 Claude Code \u2192 API Key'),
    ('1:30', '15 min', '\u9996\u6b21\u4ea4\u4e92',
     '\u6bcf\u4eba\u8dd1\u901a\u7b2c\u4e00\u6761\u547d\u4ee4'),
    ('1:45', '10 min', '\u81ea\u7531\u63a2\u7d22',
     '\u8bd5\u7528\u901f\u67e5\u5361\u547d\u4ee4'),
    ('1:55', '5 min', '\u6536\u5c3e',
     '\u603b\u7ed3 + \u9884\u544a\u7b2c 2 \u6b21\u8bfe'),
]

for i, (time, dur, phase, desc) in enumerate(timeline):
    y = 1.3 + i * 0.52
    c = ORANGE if phase in ['\u201c\u54c7\u201d\u65f6\u523b\u6f14\u793a',
                            '\u9996\u6b21\u4ea4\u4e92'] else SUBTITLE_GRAY
    if phase == '\u4f11\u606f':
        c = DIM_GRAY
    add_text(slide, Inches(0.8), Inches(y), Inches(0.9), Inches(0.4),
             time, size=14, color=DIM_GRAY, bold=True)
    add_text(slide, Inches(1.7), Inches(y), Inches(0.9), Inches(0.4),
             dur, size=12, color=DIM_GRAY)
    add_text(slide, Inches(2.8), Inches(y), Inches(3), Inches(0.4),
             phase, size=15, color=c, bold=True)
    add_text(slide, Inches(6.0), Inches(y), Inches(6.5), Inches(0.4),
             desc, size=14, color=LIGHT_GRAY)


# ============================================================
# Slide 4: Section - Opening
# ============================================================
section_slide('\u7b2c\u4e00\u90e8\u5206\uff1a\u5f00\u573a',
              '0:00 - 0:10 \xb7 \u81ea\u6211\u4ecb\u7ecd + \u94fa\u57ab')


# ============================================================
# Slide 5: Opening Script
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:00 - 0:05', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '\u5f00\u573a\u767d\uff1a\u4e00\u53e5\u8bdd\u8bf4\u6e05\u8fd9\u95e8\u8bfe',
         size=28, color=WHITE, bold=True)

lines = [
    '\u201c8 \u6b21\u8bfe\u4e4b\u540e\uff0c\u4f60\u4eec\u6bcf\u4e2a\u4eba\u90fd\u4f1a\u505a\u51fa\u4e00\u4e2a\u81ea\u5df1\u7684\u8f6f\u4ef6\u9879\u76ee\u3002\u201d',
    '\u201c\u4e0d\u9700\u8981\u4efb\u4f55\u7f16\u7a0b\u57fa\u7840\u3002\u201d',
    '\u201c\u6211\u4eec\u4e0d\u662f\u5728\u5b66\u7f16\u7a0b\u3002\u6211\u4eec\u5728\u5b66\u6307\u6325 AI \u5e2e\u4f60\u7f16\u7a0b\u3002\u201d',
    '\u201c\u4e0d\u4fe1\uff1f\u8ba9\u6211\u6f14\u793a\u7ed9\u4f60\u4eec\u770b\u3002\u201d',
]
for i, line in enumerate(lines):
    y = 1.8 + i * 1.0
    sz = 26 if i == 0 else 22
    cl = ORANGE if i == 0 else LIGHT_GRAY
    add_text(slide, Inches(1.5), Inches(y), Inches(10), Inches(0.6),
             line, size=sz, color=cl, bold=(i == 0))

add_text(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
         '\u2192 \u63a7\u5236\u5728 5 \u5206\u949f\u5185\uff0c\u7acb\u523b\u5f15\u5411\u6f14\u793a\uff0c\u4e0d\u8981\u5728\u5f00\u573a\u8bb2\u592a\u591a\u7406\u8bba',
         size=14, color=DIM_GRAY)


# ============================================================
# Slide 6: What is Claude Code
# ============================================================
teaching_slide('0:05 - 0:10',
    '\u4ec0\u4e48\u662f Claude Code\uff1f',
    [
        'Anthropic \u516c\u53f8\u505a\u7684 AI \u7f16\u7a0b\u52a9\u624b',
        '\u8dd1\u5728\u4f60\u7535\u8111\u7684\u7ec8\u7aef\u91cc',
        '\u7528\u4e2d\u6587\u544a\u8bc9\u5b83\u60f3\u505a\u4ec0\u4e48\uff0c\u5b83\u5e2e\u4f60\u5199\u4ee3\u7801\u3001\u521b\u5efa\u6587\u4ef6',
        '\u4e0d\u662f\u804a\u5929\u673a\u5668\u4eba \u2014\u2014 \u662f\u80fd\u76f4\u63a5\u5728\u4f60\u7535\u8111\u4e0a\u5e72\u6d3b\u7684 AI \u667a\u80fd\u4f53',
    ],
    note='\u2192 \u8fc7\u6e21\uff1a\u201c\u73b0\u5728\uff0c\u7ed9\u6211\u4e00\u4e2a\u60f3\u6cd5\uff0c\u6211\u6765\u6f14\u793a\u3002\u201d')


# ============================================================
# Slide 7: Section - Wow Moment
# ============================================================
section_slide('\u7b2c\u4e8c\u90e8\u5206\uff1a\u201c\u54c7\u201d\u65f6\u523b\u6f14\u793a',
              '0:10 - 0:30 \xb7 \u6574\u8282\u8bfe\u7684\u7075\u9b42')


# ============================================================
# Slide 8: Demo Step 1 - Collect Ideas
# ============================================================
teaching_slide('0:10 - 0:12',
    '\u7b2c 1 \u6b65\uff1a\u5f81\u96c6\u60f3\u6cd5',
    [
        '\u201c\u6709\u6ca1\u6709\u4eba\u60f3\u8981\u4e00\u4e2a\u7f51\u9875\uff1f\u4ec0\u4e48\u90fd\u884c\u3002\u201d',
        '\u63a5\u53d7\uff1a\u8d3a\u5361\u3001\u83dc\u5355\u3001\u4e2a\u4eba\u4e3b\u9875\u3001\u6d3b\u52a8\u9080\u8bf7\u51fd\u3001\u4ea7\u54c1\u5c55\u793a\u9875',
        '\u62d2\u7edd\uff08\u6e29\u548c\u5730\uff09\uff1a\u9700\u8981\u767b\u5f55/\u6570\u636e\u5e93/\u652f\u4ed8\u7684',
        '\u6ca1\u4eba\u8bf4\uff1f\u7528\u9884\u5907\u65b9\u6848\uff1a\u201c\u5e2e\u6211\u4eec\u73ed\u505a\u4e2a\u8bfe\u7a0b\u4ecb\u7ecd\u9875\u201d',
    ],
    note='\u2192 \u60f3\u6cd5\u5f15\u5bfc\uff1a\u5355\u9875\u3001\u9759\u6001\u3001\u89c6\u89c9\u5316\u7684\u9879\u76ee')


# ============================================================
# Slide 9: Demo Step 2 - Live Build
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:12 - 0:19', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '\u7b2c 2 \u6b65\uff1a\u73b0\u573a\u6784\u5efa\uff08\u6f14\u793a\u811a\u672c\uff09',
         size=28, color=WHITE, bold=True)

add_code(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(1.2),
         'mkdir demo && cd demo\nclaude')

add_text(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(0.4),
         '\u793a\u4f8b\u63d0\u793a\u8bcd\uff08\u5047\u8bbe\u5b66\u5458\u9009\u4e86\u201c\u9910\u5385\u83dc\u5355\u201d\uff09\uff1a',
         size=14, color=SUBTITLE_GRAY)

prompt_text = (
    '\u5e2e\u6211\u505a\u4e00\u4e2a\u4e2d\u9910\u9986\u7684\u83dc\u5355\u7f51\u9875\u3002\n'
    '\u9910\u5385\u53eb\u201c\u597d\u5473\u9053\u201d\u3002\u8981\u6709\uff1a\n'
    '- \u6f02\u4eae\u7684\u6807\u9898\u548c\u9910\u5385\u540d\u5b57\n'
    '- \u5206\u7c7b\uff1a\u51c9\u83dc\u3001\u70ed\u83dc\u3001\u6c64\u3001\u4e3b\u98df\n'
    '- \u6bcf\u4e2a\u83dc\u6709\u540d\u5b57\u548c\u4ef7\u683c\n'
    '- \u9875\u9762\u8981\u597d\u770b\uff0c\u7528\u6696\u8272\u8c03'
)
add_code(slide, Inches(0.8), Inches(3.7), Inches(5.5), Inches(2.5), prompt_text, size=13)

add_text(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4),
         '\u8fb9\u505a\u8fb9\u89e3\u8bf4\uff08\u5173\u952e\uff01\uff09\uff1a',
         size=16, color=ACCENT_BLUE, bold=True)

narration = [
    '\u201c\u770b\uff0c\u6211\u5728\u7528\u4e2d\u6587\u8ddf Claude \u8bf4\u8bdd\u201d',
    '\u201c\u5b83\u5728\u60f3... \u5728\u5206\u6790\u6211\u7684\u9700\u6c42\u201d',
    '\u201c\u5b83\u5f00\u59cb\u5199\u4ee3\u7801\u4e86\u3002\u6211\u5b8c\u5168\u4e0d\u9700\u8981\u61c2\u201d',
    '\u201c\u505a\u5b8c\u4e86\uff01\u8ba9\u6211\u4eec\u6253\u5f00\u770b\u770b\u201d',
]
add_bullets(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(3),
            narration, size=16, bullet_color=GREEN)

add_code(slide, Inches(7), Inches(5.5), Inches(5), Inches(0.8),
         'open index.html    # macOS\nstart index.html   # Windows')


# ============================================================
# Slide 10: Demo Step 3 - Show Result + Reaction
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:19 - 0:20', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '\u7b2c 3 \u6b65\uff1a\u5c55\u793a\u7ed3\u679c',
         size=28, color=WHITE, bold=True)

add_text(slide, Inches(2), Inches(2.5), Inches(9), Inches(1),
         '\u201c\u4ece\u6211\u8f93\u5165\u63cf\u8ff0\u5230\u73b0\u5728\uff0c\u5927\u6982 2 \u5206\u949f\u3002\n\u4e00\u4e2a\u5b8c\u6574\u7684\u7f51\u9875\u5c31\u505a\u597d\u4e86\u3002\u201d',
         size=26, color=WHITE)

add_text(slide, Inches(2), Inches(4.2), Inches(9), Inches(0.8),
         '\u201c8 \u6b21\u8bfe\u4e4b\u540e\uff0c\u4f60\u4eec\u6bcf\u4e2a\u4eba\u90fd\u80fd\u505a\u5230\u8fd9\u4e2a\uff0c\n\u800c\u4e14\u505a\u7684\u662f\u4f60\u81ea\u5df1\u60f3\u8981\u7684\u4e1c\u897f\u3002\u201d',
         size=24, color=ORANGE, bold=True)

add_text(slide, Inches(2), Inches(5.8), Inches(9), Inches(0.5),
         '\u7ed9 30 \u79d2\u8ba9\u5b66\u5458\u6d88\u5316\u3002\u6559\u5ba4\u5e94\u8be5\u6709\u660e\u663e\u7684\u60ca\u53f9/\u8ba8\u8bba\u58f0\u3002',
         size=15, color=DIM_GRAY)


# ============================================================
# Slide 11: Deconstruct the Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:20 - 0:30', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '\u62c6\u89e3\uff1a\u521a\u624d\u53d1\u751f\u4e86\u4ec0\u4e48\uff1f',
         size=28, color=WHITE, bold=True)

steps_data = [
    ('1', '\u63cf\u8ff0', '\u7528\u4e2d\u6587\u544a\u8bc9 Claude \u4f60\u60f3\u8981\u4ec0\u4e48'),
    ('2', '\u751f\u6210', 'Claude \u5199\u4ee3\u7801\uff0c\u521b\u5efa\u6587\u4ef6'),
    ('3', '\u8bc4\u4f30', '\u5728\u6d4f\u89c8\u5668\u91cc\u6253\u5f00\uff0c\u770b\u7ed3\u679c'),
    ('4', '\u4f18\u5316', '\u4e0d\u6ee1\u610f\u5c31\u8bf4\u51fa\u54ea\u91cc\u4e0d\u5bf9\uff0c\u8ba9 Claude \u6539'),
]
for i, (num, label, desc) in enumerate(steps_data):
    y = 1.8 + i * 1.1
    add_text(slide, Inches(1.5), Inches(y), Inches(0.8), Inches(0.6),
             num, size=36, color=ORANGE, bold=True)
    add_text(slide, Inches(2.5), Inches(y), Inches(2.5), Inches(0.5),
             label, size=24, color=WHITE, bold=True)
    add_text(slide, Inches(5.2), Inches(y + 0.05), Inches(7), Inches(0.5),
             desc, size=18, color=LIGHT_GRAY)

add_text(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
         '\u8fd9\u56db\u6b65\uff0c\u5c31\u662f\u63a5\u4e0b\u6765 8 \u6b21\u8bfe\u6211\u4eec\u8981\u7ec3\u5230\u719f\u7ec3\u7684\u4e1c\u897f\u3002',
         size=18, color=YELLOW, bold=True)


# ============================================================
# Slide 12: Course Roadmap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:30 - 0:40', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '8 \u6b21\u8bfe\u7684\u65c5\u7a0b', size=28, color=WHITE, bold=True)

roadmap = [
    ('1', '\u4eca\u5929', '\u88c5\u597d\u5de5\u5177\uff0c\u770b\u5230 AI \u7684\u5a01\u529b'),
    ('2', '\u4e0b\u6b21', '\u81ea\u5df1\u52a8\u624b\u505a\u4e1c\u897f'),
    ('3-4', '', '\u6211\u5e26\u7740\u4f60\u4eec\u505a\u4e00\u4e2a\u5b8c\u6574\u7684\u5c0f\u5de5\u5177'),
    ('5-6', '', '\u4f60\u4eec\u81ea\u5df1\u9009\u60f3\u505a\u4ec0\u4e48\uff0c\u6211\u5728\u65c1\u8fb9\u5e2e\u5fd9'),
    ('7-8', '', '\u505a\u81ea\u5df1\u7684\u9879\u76ee\uff0c\u6700\u540e\u4e00\u5929\u4e0a\u53f0\u5c55\u793a'),
]
for i, (num, label, desc) in enumerate(roadmap):
    y = 1.6 + i * 1.0
    is_today = (num == '1')
    add_text(slide, Inches(1), Inches(y), Inches(1), Inches(0.5),
             f'\u7b2c {num} \u6b21', size=18, color=ORANGE if is_today else SUBTITLE_GRAY, bold=True)
    add_text(slide, Inches(2.5), Inches(y), Inches(9), Inches(0.5),
             desc, size=20, color=WHITE if is_today else LIGHT_GRAY,
             bold=is_today)

add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         '\u201c\u5230\u7b2c 8 \u6b21\u8bfe\u7ed3\u675f\uff0c\u4f60\u624b\u91cc\u4f1a\u6709\u4e00\u4e2a\u4f60\u81ea\u5df1\u505a\u7684\u3001\u53ef\u4ee5\u7528\u7684\u8f6f\u4ef6\u4f5c\u54c1\u3002\u201d',
         size=17, color=ORANGE)


# ============================================================
# Slide 13: Section - Terminal
# ============================================================
section_slide('\u7b2c\u4e09\u90e8\u5206\uff1a\u7ec8\u7aef\u5165\u95e8',
              '0:40 - 0:55 \xb7 \u8ba4\u8bc6\u4f60\u7684\u65b0\u5de5\u5177')


# ============================================================
# Slide 14: Open Terminal
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:40', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '\u7ec8\u7aef\u662f\u4ec0\u4e48\uff1f',
         size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.6),
         '\u5e73\u65f6\u4f60\u7528\u9f20\u6807\u70b9\u6765\u70b9\u53bb\uff0c\u7ec8\u7aef\u662f\u4f60\u7528\u6253\u5b57\u7684\u65b9\u5f0f\u544a\u8bc9\u7535\u8111\u505a\u4ec0\u4e48\u3002',
         size=20, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(0.5),
         'macOS \u5b66\u5458', size=20, color=ACCENT_BLUE, bold=True)
add_code(slide, Inches(0.8), Inches(3.4), Inches(5.5), Inches(0.8),
         'Command + \u7a7a\u683c \u2192 \u8f93\u5165 Terminal \u2192 \u56de\u8f66')

add_text(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(0.5),
         'Windows \u5b66\u5458', size=20, color=ACCENT_BLUE, bold=True)
add_code(slide, Inches(7), Inches(3.4), Inches(5.5), Inches(0.8),
         'Windows \u952e \u2192 \u8f93\u5165 PowerShell \u2192 \u70b9\u51fb\u6253\u5f00')

add_text(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
         '\u52a9\u6559\u6b64\u65f6\u8d70\u52a8\uff0c\u786e\u4fdd\u6bcf\u4eba\u90fd\u6253\u5f00\u4e86\u7ec8\u7aef',
         size=15, color=YELLOW, bold=True)


# ============================================================
# Slide 15: Basic Commands
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '0:45 - 0:55', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '5 \u4e2a\u57fa\u7840\u547d\u4ee4 \u2014\u2014 \u8fd9\u5c31\u591f\u4e86',
         size=28, color=WHITE, bold=True)

cmds = [
    ('pwd', '\u663e\u793a\u5f53\u524d\u4f4d\u7f6e',
     '\u201c\u544a\u8bc9\u4f60\u73b0\u5728\u5728\u7535\u8111\u7684\u54ea\u4e2a\u6587\u4ef6\u5939\u91cc\u201d'),
    ('ls', '\u5217\u51fa\u6587\u4ef6',
     '\u201c\u5217\u51fa\u5f53\u524d\u6587\u4ef6\u5939\u91cc\u6709\u4ec0\u4e48\u201d\uff08Win: dir\uff09'),
    ('mkdir test', '\u521b\u5efa\u6587\u4ef6\u5939',
     '\u201c\u521b\u5efa\u4e86\u4e00\u4e2a\u53eb test \u7684\u65b0\u6587\u4ef6\u5939\u201d'),
    ('cd test', '\u8fdb\u5165\u6587\u4ef6\u5939',
     '\u201c\u8fdb\u5165\u4e86 test \u6587\u4ef6\u5939\u201d'),
    ('cd ..', '\u56de\u5230\u4e0a\u4e00\u5c42',
     '\u201c\u56de\u5230\u4e0a\u4e00\u5c42\u201d'),
]
for i, (cmd, action, say) in enumerate(cmds):
    y = 1.6 + i * 1.05
    add_code(slide, Inches(0.8), Inches(y), Inches(2.5), Inches(0.5), cmd, size=18)
    add_text(slide, Inches(3.5), Inches(y), Inches(2.5), Inches(0.4),
             action, size=16, color=ORANGE, bold=True)
    add_text(slide, Inches(6.2), Inches(y), Inches(6.5), Inches(0.4),
             say, size=15, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
         '\u68c0\u67e5\u70b9\uff1a\u201c\u8bf7\u6bcf\u4e2a\u4eba\u8dd1\u4e00\u4e0b pwd\uff0c\u770b\u5230\u7ed3\u679c\u7684\u4e3e\u624b\u3002\u201d',
         size=15, color=YELLOW, bold=True)


# ============================================================
# Slide 16: Break
# ============================================================
section_slide('\u2615 \u4f11\u606f 10 \u5206\u949f',
              '\u7ec8\u7aef\u64cd\u4f5c\u9047\u5230\u95ee\u9898\uff1f\u4f11\u606f\u65f6\u627e\u52a9\u6559\u5e2e\u5fd9')


# ============================================================
# Slide 17: Section - Install
# ============================================================
section_slide('\u7b2c\u56db\u90e8\u5206\uff1a\u73af\u5883\u5b89\u88c5',
              '1:05 - 1:30 \xb7 \u628a Claude Code \u88c5\u5230\u4f60\u7684\u7535\u8111\u4e0a')


# ============================================================
# Slide 18: Install Node.js
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '\u6b65\u9aa4 1/3', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '\u5b89\u88c5 Node.js', size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
         'macOS', size=18, color=ACCENT_BLUE, bold=True)
add_code(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(1.5),
         '# \u68c0\u67e5\u662f\u5426\u5df2\u5b89\u88c5\nnode --version\n\n# \u6ca1\u6709\uff1f\u7528 USB \u4e2d\u7684\u5b89\u88c5\u5305\n# \u53cc\u51fb node-vXX.X.X.pkg \u5b89\u88c5\n\n# \u91cd\u65b0\u6253\u5f00\u7ec8\u7aef\uff0c\u9a8c\u8bc1\nnode --version')

add_text(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
         'Windows', size=18, color=ACCENT_BLUE, bold=True)
add_code(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(1.5),
         '# \u68c0\u67e5\u662f\u5426\u5df2\u5b89\u88c5\nnode --version\n\n# \u6ca1\u6709\uff1f\u7528 USB \u4e2d\u7684\u5b89\u88c5\u5305\n# \u53cc\u51fb node-vXX.X.X.msi\n# \u52fe\u9009 "Add to PATH"\n\n# \u91cd\u65b0\u6253\u5f00 PowerShell\uff0c\u9a8c\u8bc1\nnode --version')

add_text(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
         '\u6210\u529f\u6807\u5fd7\uff1a\u770b\u5230\u7248\u672c\u53f7\uff08\u5982 v22.x.x\uff09\u2192 \u201c\u770b\u5230\u7684\u540c\u5b66\u4e3e\u624b\uff01\u201d',
         size=16, color=YELLOW, bold=True)

# Troubleshooting
add_text(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.4),
         '\u5e38\u89c1\u95ee\u9898\uff1acommand not found \u2192 \u91cd\u5f00\u7ec8\u7aef\uff5c\u65e0\u6cd5\u6253\u5f00\u5b89\u88c5\u5305 \u2192 \u7cfb\u7edf\u5b89\u5168\u8bbe\u7f6e\u5141\u8bb8\uff5c\u65e0\u6743\u9650 \u2192 \u7528\u7f51\u9875\u7248\u5907\u7528',
         size=13, color=DIM_GRAY)


# ============================================================
# Slide 19: Install Claude Code
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '\u6b65\u9aa4 2/3', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '\u5b89\u88c5 Claude Code', size=28, color=WHITE, bold=True)

add_code(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.2),
         '# \u5b89\u88c5\nnpm install -g @anthropic-ai/claude-code\n\n# \u9a8c\u8bc1\nclaude --version')

add_text(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.5),
         '\u6210\u529f\u6807\u5fd7\uff1a\u770b\u5230\u7248\u672c\u53f7 \u2192 \u201c\u4e3e\u624b\uff01\u201d',
         size=16, color=YELLOW, bold=True)

add_text(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.4),
         '\u5e38\u89c1\u95ee\u9898\uff1a', size=16, color=SUBTITLE_GRAY, bold=True)
issues = [
    'EACCES \u6743\u9650\u9519\u8bef \u2192 macOS: sudo npm install -g ...  Windows: \u7ba1\u7406\u5458\u8eab\u4efd\u8fd0\u884c',
    '\u7f51\u7edc\u8d85\u65f6 \u2192 \u7528 USB \u79bb\u7ebf\u5305: npm install -g ./claude-code-x.x.x.tgz',
    'npm: command not found \u2192 \u56de\u5230\u6b65\u9aa4 1\uff0c\u91cd\u65b0\u5b89\u88c5 Node.js',
]
add_bullets(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(2),
            issues, size=14, color=LIGHT_GRAY, bullet_color=RED_SOFT)


# ============================================================
# Slide 20: Configure API Key
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '\u6b65\u9aa4 3/3', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '\u914d\u7f6e API Key', size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
         '\u201cClaude Code \u9700\u8981\u4e00\u628a\u2018\u94a5\u5319\u2019\u624d\u80fd\u5de5\u4f5c\u3002\u6211\u5df2\u7ecf\u4e3a\u5927\u5bb6\u51c6\u5907\u597d\u4e86\u3002\u201d',
         size=20, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(0.4),
         'macOS / Linux', size=18, color=ACCENT_BLUE, bold=True)
add_code(slide, Inches(0.8), Inches(3.1), Inches(5.5), Inches(0.7),
         'export ANTHROPIC_API_KEY="sk-ant-xxxxx"')

add_text(slide, Inches(7), Inches(2.6), Inches(5.5), Inches(0.4),
         'Windows PowerShell', size=18, color=ACCENT_BLUE, bold=True)
add_code(slide, Inches(7), Inches(3.1), Inches(5.5), Inches(0.7),
         '$env:ANTHROPIC_API_KEY="sk-ant-xxxxx"')

add_text(slide, Inches(0.8), Inches(4.3), Inches(11), Inches(0.4),
         '\u5206\u7ec4\u5206\u914d\uff1a', size=16, color=SUBTITLE_GRAY, bold=True)
groups = [
    '1-5 \u53f7\u5b66\u5458\uff1aKey A',
    '6-10 \u53f7\u5b66\u5458\uff1aKey B',
    '11-15 \u53f7\u5b66\u5458\uff1aKey C',
]
add_bullets(slide, Inches(0.8), Inches(4.8), Inches(5), Inches(1.5),
            groups, size=16, bullet_color=ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
         '\u8fd9\u4e2a Key \u662f\u8bfe\u5802\u4e13\u7528\u7684\uff0c\u8bf7\u4e0d\u8981\u5206\u4eab\u7ed9\u522b\u4eba\u3002\u8bfe\u7a0b\u7ed3\u675f\u540e\u4f1a\u6559\u4f60\u4eec\u7533\u8bf7\u81ea\u5df1\u7684\u3002',
         size=14, color=DIM_GRAY)


# ============================================================
# Slide 21: Section - First Interaction
# ============================================================
section_slide('\u7b2c\u4e94\u90e8\u5206\uff1a\u9996\u6b21\u4ea4\u4e92',
              '1:30 - 1:45 \xb7 \u6bcf\u4e2a\u4eba\u90fd\u8981\u8dd1\u901a\u4e00\u6b21\uff01')


# ============================================================
# Slide 22: First Command
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '1:30', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '\u4f60\u7684\u7b2c\u4e00\u4e2a\u9879\u76ee\uff01',
         size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.4),
         '\u8ddf\u6211\u4e00\u8d77\u505a\uff1a', size=18, color=SUBTITLE_GRAY)

add_code(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(2.5),
         '# 1. \u521b\u5efa\u7ec3\u4e60\u6587\u4ef6\u5939\n'
         'mkdir lesson1 && cd lesson1\n\n'
         '# 2. \u542f\u52a8 Claude Code\n'
         'claude\n\n'
         '# 3. \u8f93\u5165\u4f60\u7684\u7b2c\u4e00\u6761\u6307\u4ee4\n'
         '\u4f60\u597d\uff01\u8bf7\u5e2e\u6211\u521b\u5efa\u4e00\u4e2a\u7b80\u5355\u7684 HTML \u9875\u9762\uff0c\n'
         '\u5185\u5bb9\u662f\u201c\u6211\u7684\u7b2c\u4e00\u4e2a\u7f51\u9875\u201d\uff0c\u6807\u9898\u7528\u5927\u5b57\uff0c\n'
         '\u9875\u9762\u80cc\u666f\u7528\u6d45\u84dd\u8272\u3002', size=15)

add_text(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.4),
         '\u7b49 Claude \u5b8c\u6210\u540e\uff1a', size=16, color=SUBTITLE_GRAY)
add_code(slide, Inches(0.8), Inches(5.7), Inches(5), Inches(0.5),
         'open index.html    # macOS', size=15)
add_code(slide, Inches(6.5), Inches(5.7), Inches(5), Inches(0.5),
         'start index.html   # Windows', size=15)

add_text(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
         '\u201c\u770b\u5230\u7f51\u9875\u7684\u540c\u5b66\u4e3e\u624b\uff01\u201d \u2014\u2014 \u76ee\u6807\uff1a95%+ \u4e3e\u624b',
         size=16, color=YELLOW, bold=True)


# ============================================================
# Slide 23: Free Explore
# ============================================================
teaching_slide('1:45 - 1:55',
    '\u81ea\u7531\u63a2\u7d22\uff1a\u8bd5\u8bd5\u901f\u67e5\u5361',
    [
        '\u201c\u8bf7\u5e2e\u6211\u4fee\u6539\u7f51\u9875\u7684\u80cc\u666f\u8272\u4e3a\u4f60\u559c\u6b22\u7684\u989c\u8272\u201d',
        '\u201c\u7ed9\u7f51\u9875\u52a0\u4e00\u5f20\u56fe\u7247\u63cf\u8ff0\uff08\u7528 emoji \u4ee3\u66ff\uff09\u201d',
        '\u201c\u8ba9\u8fd9\u4e2a\u9875\u9762\u5728\u624b\u673a\u4e0a\u4e5f\u597d\u770b\u201d',
        '\u901f\u67e5\u5361\u4e0a\u8fd8\u6709\u66f4\u591a\u547d\u4ee4\uff0c\u968f\u4fbf\u8bd5\uff01',
    ],
    note='\u52a9\u6559\uff1a\u8d70\u52a8\u89c2\u5bdf\uff0c\u770b\u5230\u6709\u8da3\u7684\u6210\u679c\u8ba9\u5b66\u5458\u5c55\u793a\u7ed9\u65c1\u8fb9\u7684\u4eba\u770b')


# ============================================================
# Slide 24: Wrap Up
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.4),
         '1:55 - 2:00', size=13, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '\u4eca\u5929\u5b8c\u6210\u4e86\u4e09\u4ef6\u4e8b',
         size=28, color=WHITE, bold=True)

achievements = [
    ('1', '\u4f60\u770b\u5230\u4e86 AI \u80fd\u505a\u4ec0\u4e48',
     '10 \u5206\u949f\u505a\u51fa\u4e00\u4e2a\u7f51\u9875'),
    ('2', '\u4f60\u628a\u5de5\u5177\u88c5\u597d\u4e86',
     'Claude Code \u5df2\u5c31\u7eea'),
    ('3', '\u4f60\u81ea\u5df1\u8dd1\u901a\u4e86\u7b2c\u4e00\u4e2a\u9879\u76ee',
     '\u201c\u6211\u7684\u7b2c\u4e00\u4e2a\u7f51\u9875\u201d'),
]
for i, (num, title, desc) in enumerate(achievements):
    y = 1.8 + i * 1.3
    add_text(slide, Inches(1.5), Inches(y), Inches(0.6), Inches(0.5),
             num, size=32, color=ORANGE, bold=True)
    add_text(slide, Inches(2.3), Inches(y), Inches(9), Inches(0.5),
             title, size=22, color=WHITE, bold=True)
    add_text(slide, Inches(2.3), Inches(y + 0.5), Inches(9), Inches(0.4),
             desc, size=16, color=SUBTITLE_GRAY)


# ============================================================
# Slide 25: Homework + Preview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         '\u8bfe\u540e\u4efb\u52a1\uff08\u53ef\u9009\uff0c\u4e0d\u5f3a\u5236\uff09',
         size=28, color=WHITE, bold=True)

add_bullets(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(2.5),
    [
        '\u8bd5\u7740\u7528 Claude Code \u505a\u4e00\u4e2a\u4f60\u81ea\u5df1\u60f3\u8981\u7684\u7b80\u5355\u7f51\u9875',
        '\u4efb\u4f55\u4e1c\u897f\u90fd\u884c\uff1a\u4e00\u5f20\u8d3a\u5361\u3001\u4e00\u4e2a\u5012\u8ba1\u65f6\u3001\u4e00\u4e2a\u4e2a\u4eba\u4ecb\u7ecd',
        '\u505a\u4e0d\u51fa\u6765\u4e5f\u5b8c\u5168\u6ca1\u5173\u7cfb\uff0c\u4e0b\u6b21\u8bfe\u6211\u4eec\u4e00\u8d77\u6765',
    ], size=18)

add_text(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.7),
         '\u4e0b\u6b21\u8bfe\u9884\u544a', size=24, color=ORANGE, bold=True)

add_text(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(0.8),
         '\u7b2c 2 \u6b21\u8bfe\uff1aAI \u601d\u7ef4\u6a21\u578b + \u521d\u6b21\u52a8\u624b\n'
         '\u5b66\u4e60\u5982\u4f55\u66f4\u597d\u5730\u8ddf Claude \u5bf9\u8bdd\uff0c\u8ba9\u5b83\u505a\u51fa\u4f60\u771f\u6b63\u60f3\u8981\u7684\u4e1c\u897f',
         size=18, color=LIGHT_GRAY)


# ============================================================
# Slide 26: Teacher Notes - Key Reminders
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         '\u8001\u5e08\u5907\u6ce8\uff1a5 \u6761\u5173\u952e\u63d0\u9192',
         size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
         '\uff08\u6b64\u9875\u4e0d\u5c55\u793a\u7ed9\u5b66\u5458\uff0c\u4ec5\u8001\u5e08\u53c2\u8003\uff09',
         size=13, color=DIM_GRAY)

reminders = [
    ('\u201c\u54c7\u201d\u65f6\u523b\u662f\u6574\u8282\u8bfe\u7684\u7075\u9b42',
     '\u5982\u679c\u6f14\u793a\u4e0d\u9707\u64bc\uff0c\u540e\u9762\u7684\u5b89\u88c5\u73af\u8282\u4f1a\u53d8\u6210\u714e\u71ac\u3002\u82b1\u65f6\u95f4\u51c6\u5907\u548c\u7ec3\u4e60\u6f14\u793a\u3002'),
    ('\u5b89\u88c5\u95ee\u9898\u4e0d\u8981\u604b\u6218',
     '\u4e00\u4e2a\u5b66\u5458\u5361\u8d85\u8fc7 5 \u5206\u949f\u5c31\u5207\u7f51\u9875\u7248\u3002\u4fdd\u62a4\u6574\u4f53\u8282\u594f\u6bd4\u89e3\u51b3\u4e2a\u522b\u95ee\u9898\u91cd\u8981\u3002'),
    ('\u8ba9\u5b66\u5458\u52a8\u624b\uff0c\u4e0d\u8981\u53ea\u662f\u770b',
     '\u9996\u6b21\u4ea4\u4e92\u73af\u8282\u5fc5\u987b\u6bcf\u4eba\u90fd\u505a\u4e00\u6b21\u3002\u624b\u6307\u6572\u8fc7\u952e\u76d8\u548c\u773c\u775b\u770b\u8fc7\u5c4f\u5e55\uff0c\u8bb0\u5fc6\u6548\u679c\u5b8c\u5168\u4e0d\u540c\u3002'),
    ('\u8425\u9020\u5b89\u5168\u611f',
     '\u53cd\u590d\u8bf4\uff1a\u201c\u505a\u4e0d\u51fa\u6765\u5b8c\u5168\u6b63\u5e38\u201d\u201c\u6ca1\u6709\u8822\u95ee\u9898\u201d\u3002\u96f6\u57fa\u7840\u7684\u6210\u5e74\u4eba\u6700\u6015\u5728\u522b\u4eba\u9762\u524d\u663e\u5f97\u7b28\u3002'),
    ('\u4e0d\u8981\u63d0\u524d\u6559\u592a\u591a\u6982\u5ff5',
     '\u4eca\u5929\u7684\u76ee\u6807\u53ea\u6709\u4e24\u4e2a\uff1a\u9707\u64bc + \u5b89\u88c5\u3002\u601d\u7ef4\u6a21\u578b\u7559\u7ed9\u7b2c 2 \u6b21\u8bfe\u3002'),
]
for i, (title, desc) in enumerate(reminders):
    y = 1.7 + i * 0.95
    add_text(slide, Inches(0.8), Inches(y), Inches(0.5), Inches(0.4),
             str(i + 1), size=18, color=ORANGE, bold=True)
    add_text(slide, Inches(1.4), Inches(y), Inches(4), Inches(0.4),
             title, size=16, color=WHITE, bold=True)
    add_text(slide, Inches(5.5), Inches(y), Inches(7), Inches(0.4),
             desc, size=14, color=LIGHT_GRAY)


# ============================================================
# Slide 27: Backup Plans
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         '\u8001\u5e08\u5907\u6ce8\uff1a\u5907\u7528\u65b9\u6848',
         size=28, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
         '\uff08\u6b64\u9875\u4e0d\u5c55\u793a\u7ed9\u5b66\u5458\uff0c\u4ec5\u8001\u5e08\u53c2\u8003\uff09',
         size=13, color=DIM_GRAY)

backups = [
    ('API \u6545\u969c', '\u64ad\u653e\u9884\u5f55\u5c4f\u5e55\u5f55\u50cf'),
    ('Wi-Fi \u5d29\u6e83', 'USB \u79bb\u7ebf\u5b89\u88c5\u5305'),
    ('\u7535\u8111\u65e0\u6cd5\u5b89\u88c5', 'claude.ai/code \u7f51\u9875\u7248'),
    ('\u60f3\u6cd5\u592a\u590d\u6742', '\u6e29\u548c\u5f15\u5bfc\u5230\u9884\u5907\u65b9\u6848'),
    ('\u5b89\u88c5\u65f6\u95f4\u4e0d\u591f', '\u7b2c 2 \u6b21\u8bfe\u524d 20 \u5206\u949f\u7ee7\u7eed'),
    ('\u5b66\u5458\u5b8c\u5168\u8ddf\u4e0d\u4e0a', '\u52a9\u6559\u4e00\u5bf9\u4e00\uff0c\u5148\u7528\u7f51\u9875\u7248'),
]
for i, (scenario, solution) in enumerate(backups):
    y = 1.8 + i * 0.85
    add_text(slide, Inches(0.8), Inches(y), Inches(3), Inches(0.4),
             scenario, size=17, color=RED_SOFT, bold=True)
    add_text(slide, Inches(4.2), Inches(y), Inches(8.5), Inches(0.4),
             solution, size=16, color=LIGHT_GRAY)


# ============================================================
# Save
# ============================================================
output = '/Users/wenfengzhu/ClaudeProjects/ClaudeClass/\u7b2c1\u6b21\u8bfe-\u6559\u5b66PPT.pptx'
prs.save(output)
print(f'PPT saved to: {output}')
